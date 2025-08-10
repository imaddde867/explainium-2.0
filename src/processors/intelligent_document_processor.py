"""
EXPLAINIUM - Intelligent Document Processor

Advanced multi-modal document processing with deep understanding,
optimized for Apple M4 Mac. Supports all document types with
sophisticated knowledge extraction and contextual analysis.
"""

import os
import io
import json
import logging
import asyncio
from typing import Dict, Any, List, Optional, Tuple, Union
from pathlib import Path
from datetime import datetime
import tempfile
import mimetypes

# Core processing libraries
import requests
import cv2
import numpy as np
from PIL import Image, ImageEnhance
import pandas as pd
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

# Advanced AI libraries
import torch
from transformers import pipeline
import whisper
from sentence_transformers import SentenceTransformer

# OCR and Vision
import pytesseract
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False

# Audio processing
try:
    from pyannote.audio import Pipeline as PyannoteAudioPipeline
    PYANNOTE_AVAILABLE = True
except ImportError:
    PYANNOTE_AVAILABLE = False

# Internal imports
from src.logging_config import get_logger
from src.core.config import config_manager
from src.core.optimization import ModelManager, BatchProcessor, StreamingProcessor, setup_m4_optimization
from src.ai.advanced_knowledge_engine import AdvancedKnowledgeEngine, KnowledgeEntity
from src.exceptions import ProcessingError, AIError

logger = get_logger(__name__)

# Setup M4 optimizations
setup_m4_optimization()


class MultiModalProcessor:
    """Multi-modal content processor with AI understanding"""
    
    def __init__(self):
        self.model_manager = ModelManager()
        self.knowledge_engine = AdvancedKnowledgeEngine()
        
        # Initialize OCR engines
        self._init_ocr()
        
        # Initialize audio processing
        self._init_audio()
        
        # Initialize vision models
        self._init_vision()
    
    def _init_ocr(self):
        """Initialize OCR capabilities"""
        try:
            # Test pytesseract
            pytesseract.get_tesseract_version()
            self.tesseract_available = True
            logger.info("Tesseract OCR initialized")
            
            # Initialize EasyOCR if available
            if EASYOCR_AVAILABLE:
                self.easyocr_reader = easyocr.Reader(['en'])
                self.easyocr_available = True
                logger.info("EasyOCR initialized")
            else:
                self.easyocr_available = False
                
        except Exception as e:
            logger.warning(f"OCR initialization failed: {e}")
            self.tesseract_available = False
            self.easyocr_available = False
    
    def _init_audio(self):
        """Initialize audio processing"""
        try:
            # Initialize Whisper
            self.whisper_model = whisper.load_model("base")
            self.whisper_available = True
            logger.info("Whisper model loaded")
            
            # Initialize speaker diarization if available
            if PYANNOTE_AVAILABLE:
                try:
                    self.diarization_pipeline = PyannoteAudioPipeline.from_pretrained(
                        "pyannote/speaker-diarization@2.1",
                        use_auth_token=False
                    )
                    self.diarization_available = True
                    logger.info("Speaker diarization initialized")
                except Exception as e:
                    logger.warning(f"Speaker diarization failed to initialize: {e}")
                    self.diarization_available = False
            else:
                self.diarization_available = False
                
        except Exception as e:
            logger.warning(f"Audio processing initialization failed: {e}")
            self.whisper_available = False
            self.diarization_available = False
    
    def _init_vision(self):
        """Initialize vision models"""
        try:
            # Document layout understanding
            self.layout_processor = pipeline(
                "document-question-answering",
                model="microsoft/layoutlmv3-base",
                device="mps" if torch.backends.mps.is_available() else "cpu"
            )
            
            # Image captioning
            self.image_captioner = pipeline(
                "image-to-text",
                model="Salesforce/blip-image-captioning-base",
                device="mps" if torch.backends.mps.is_available() else "cpu"
            )
            
            self.vision_available = True
            logger.info("Vision models initialized")
            
        except Exception as e:
            logger.warning(f"Vision models initialization failed: {e}")
            self.vision_available = False
    
    async def process_image(self, image_path: str, context: str = "") -> Dict[str, Any]:
        """Process image with OCR and vision understanding"""
        try:
            image = Image.open(image_path).convert('RGB')
            
            # Extract text using multiple OCR methods
            ocr_results = await self._extract_text_from_image(image)
            
            # Generate image description
            description = ""
            if self.vision_available:
                try:
                    caption_result = self.image_captioner(image)
                    description = caption_result[0]['generated_text']
                except Exception as e:
                    logger.error(f"Error generating image caption: {e}")
            
            # Combine OCR text and description
            combined_text = f"{description}\n\n{ocr_results.get('text', '')}"
            
            return {
                'text': combined_text,
                'ocr_results': ocr_results,
                'description': description,
                'image_info': {
                    'size': image.size,
                    'mode': image.mode,
                    'format': image.format
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing image: {e}")
            return {'text': '', 'error': str(e)}
    
    async def _extract_text_from_image(self, image: Image.Image) -> Dict[str, Any]:
        """Extract text using multiple OCR methods"""
        results = {'text': '', 'confidence': 0.0, 'method': 'none'}
        
        try:
            # Try EasyOCR first (usually better quality)
            if self.easyocr_available:
                try:
                    image_array = np.array(image)
                    easyocr_results = self.easyocr_reader.readtext(image_array)
                    
                    text_parts = []
                    confidences = []
                    
                    for (bbox, text, conf) in easyocr_results:
                        if conf > 0.5:  # Only include high-confidence text
                            text_parts.append(text)
                            confidences.append(conf)
                    
                    if text_parts:
                        results['text'] = ' '.join(text_parts)
                        results['confidence'] = sum(confidences) / len(confidences)
                        results['method'] = 'easyocr'
                        return results
                        
                except Exception as e:
                    logger.warning(f"EasyOCR failed: {e}")
            
            # Fallback to Tesseract
            if self.tesseract_available:
                try:
                    # Enhance image for better OCR
                    enhanced_image = self._enhance_image_for_ocr(image)
                    
                    # Extract text with confidence
                    data = pytesseract.image_to_data(enhanced_image, output_type=pytesseract.Output.DICT)
                    
                    text_parts = []
                    confidences = []
                    
                    for i, conf in enumerate(data['conf']):
                        if int(conf) > 30:  # Tesseract confidence threshold
                            word = data['text'][i].strip()
                            if word:
                                text_parts.append(word)
                                confidences.append(int(conf))
                    
                    if text_parts:
                        results['text'] = ' '.join(text_parts)
                        results['confidence'] = sum(confidences) / len(confidences) / 100  # Normalize to 0-1
                        results['method'] = 'tesseract'
                        
                except Exception as e:
                    logger.warning(f"Tesseract OCR failed: {e}")
            
        except Exception as e:
            logger.error(f"OCR extraction failed: {e}")
        
        return results
    
    def _enhance_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Enhance image quality for better OCR results"""
        try:
            # Convert to grayscale
            gray = image.convert('L')
            
            # Enhance contrast
            enhancer = ImageEnhance.Contrast(gray)
            enhanced = enhancer.enhance(2.0)
            
            # Enhance sharpness
            sharpness_enhancer = ImageEnhance.Sharpness(enhanced)
            enhanced = sharpness_enhancer.enhance(1.5)
            
            return enhanced
            
        except Exception as e:
            logger.warning(f"Image enhancement failed: {e}")
            return image
    
    async def process_audio(self, audio_path: str) -> Dict[str, Any]:
        """Process audio with transcription and speaker diarization"""
        try:
            if not self.whisper_available:
                return {'text': '', 'error': 'Whisper not available'}
            
            # Transcribe audio
            result = self.whisper_model.transcribe(audio_path)
            transcript = result['text']
            
            # Add speaker diarization if available
            speakers = []
            if self.diarization_available:
                try:
                    diarization = self.diarization_pipeline(audio_path)
                    
                    # Extract speaker segments
                    for turn, _, speaker in diarization.itertracks(yield_label=True):
                        speakers.append({
                            'speaker': speaker,
                            'start': turn.start,
                            'end': turn.end
                        })
                        
                except Exception as e:
                    logger.warning(f"Speaker diarization failed: {e}")
            
            return {
                'text': transcript,
                'language': result.get('language', 'unknown'),
                'speakers': speakers,
                'segments': result.get('segments', [])
            }
            
        except Exception as e:
            logger.error(f"Audio processing failed: {e}")
            return {'text': '', 'error': str(e)}
    
    async def process_video(self, video_path: str) -> Dict[str, Any]:
        """Process video with audio extraction and frame analysis"""
        try:
            # Extract audio from video
            audio_path = video_path.replace(Path(video_path).suffix, '_temp.wav')
            
            # Use ffmpeg to extract audio
            import subprocess
            cmd = [
                'ffmpeg', '-i', video_path, '-vn', '-acodec', 'pcm_s16le',
                '-ar', '16000', '-ac', '1', audio_path, '-y'
            ]
            
            subprocess.run(cmd, check=True, capture_output=True)
            
            # Process extracted audio
            audio_results = await self.process_audio(audio_path)
            
            # Clean up temporary audio file
            if os.path.exists(audio_path):
                os.remove(audio_path)
            
            # TODO: Add frame analysis for visual content
            
            return {
                'text': audio_results.get('text', ''),
                'audio_analysis': audio_results,
                'video_info': self._get_video_info(video_path)
            }
            
        except Exception as e:
            logger.error(f"Video processing failed: {e}")
            return {'text': '', 'error': str(e)}
    
    def _get_video_info(self, video_path: str) -> Dict[str, Any]:
        """Get basic video information"""
        try:
            import cv2
            cap = cv2.VideoCapture(video_path)
            
            info = {
                'duration': cap.get(cv2.CAP_PROP_FRAME_COUNT) / cap.get(cv2.CAP_PROP_FPS),
                'fps': cap.get(cv2.CAP_PROP_FPS),
                'width': int(cap.get(cv2.CAP_PROP_FRAME_WIDTH)),
                'height': int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT)),
                'frame_count': int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            }
            
            cap.release()
            return info
            
        except Exception as e:
            logger.warning(f"Failed to get video info: {e}")
            return {}


class IntelligentDocumentProcessor:
    """
    Advanced document processor with multi-modal understanding and
    deep knowledge extraction capabilities.
    """
    
    def __init__(self):
        self.config = config_manager.get_ai_config()
        self.knowledge_engine = AdvancedKnowledgeEngine()
        self.multimodal_processor = MultiModalProcessor()
        self.batch_processor = BatchProcessor(batch_size=self.config.batch_size)
        self.streaming_processor = StreamingProcessor(
            chunk_size=self.config.chunk_size,
            overlap=self.config.chunk_overlap
        )
        
        # Supported formats with enhanced capabilities
        self.supported_formats = {
            'text': ['.pdf', '.doc', '.docx', '.txt', '.rtf', '.md'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'],
            'spreadsheet': ['.xls', '.xlsx', '.csv'],
            'presentation': ['.ppt', '.pptx'],
            'audio': ['.mp3', '.wav', '.flac', '.aac', '.m4a'],
            'video': ['.mp4', '.avi', '.mov', '.mkv', '.webm']
        }
        
        logger.info("IntelligentDocumentProcessor initialized")
    
    async def process_document(self, file_path: str, document_id: int, 
                             company_context: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        Main method to process any supported document with deep understanding
        """
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                raise ProcessingError(f"File not found: {file_path}")
            
            # Determine file type and processing strategy
            file_type = self._determine_file_type(file_path)
            
            logger.info(f"Processing document {document_id}: {file_path.name} as {file_type}")
            
            # Extract content based on file type
            extraction_result = await self._extract_content(file_path, file_type)
            
            if not extraction_result.get('text'):
                logger.warning(f"No text content extracted from {file_path}")
                return {'error': 'No extractable content found'}
            
            # Prepare document for knowledge extraction
            document = {
                'id': str(document_id),
                'content': extraction_result['text'],
                'metadata': {
                    'filename': file_path.name,
                    'file_type': file_type,
                    'extraction_method': extraction_result.get('method', 'unknown'),
                    'confidence': extraction_result.get('confidence', 1.0),
                    'company_context': company_context or {}
                }
            }
            
            # Perform deep knowledge extraction
            knowledge_results = await self.knowledge_engine.extract_deep_knowledge(document)
            
            # Build knowledge graph
            knowledge_graph = await self.knowledge_engine.build_knowledge_graph(knowledge_results)
            
            # Extract operational intelligence
            operational_intelligence = await self.knowledge_engine.extract_operational_intelligence(
                extraction_result['text']
            )
            
            # Compile final results
            results = {
                'document_id': document_id,
                'filename': file_path.name,
                'file_type': file_type,
                'extraction_metadata': extraction_result,
                'knowledge_extraction': knowledge_results,
                'operational_intelligence': operational_intelligence,
                'knowledge_graph_summary': knowledge_graph.get_knowledge_summary(),
                'processing_timestamp': datetime.now().isoformat(),
                'confidence_scores': self._calculate_confidence_scores(knowledge_results)
            }
            
            logger.info(f"Document {document_id} processed successfully")
            return results
            
        except Exception as e:
            logger.error(f"Error processing document {document_id}: {e}")
            return {
                'document_id': document_id,
                'error': str(e),
                'processing_timestamp': datetime.now().isoformat()
            }
    
    def _determine_file_type(self, file_path: Path) -> str:
        """Determine file type and processing strategy"""
        extension = file_path.suffix.lower()
        
        for category, extensions in self.supported_formats.items():
            if extension in extensions:
                return category
        
        # Try to determine by MIME type
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type:
            if mime_type.startswith('text/'):
                return 'text'
            elif mime_type.startswith('image/'):
                return 'image'
            elif mime_type.startswith('audio/'):
                return 'audio'
            elif mime_type.startswith('video/'):
                return 'video'
        
        # Default to text processing
        return 'text'
    
    async def _extract_content(self, file_path: Path, file_type: str) -> Dict[str, Any]:
        """Extract content based on file type"""
        try:
            if file_type == 'text':
                return await self._extract_text_content(file_path)
            elif file_type == 'image':
                return await self.multimodal_processor.process_image(str(file_path))
            elif file_type == 'audio':
                return await self.multimodal_processor.process_audio(str(file_path))
            elif file_type == 'video':
                return await self.multimodal_processor.process_video(str(file_path))
            elif file_type == 'spreadsheet':
                return await self._extract_spreadsheet_content(file_path)
            elif file_type == 'presentation':
                return await self._extract_presentation_content(file_path)
            else:
                return await self._extract_text_content(file_path)
                
        except Exception as e:
            logger.error(f"Content extraction failed for {file_path}: {e}")
            return {'text': '', 'error': str(e)}
    
    async def _extract_text_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from text-based documents"""
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.pdf':
                return await self._extract_pdf_content(file_path)
            elif extension in ['.doc', '.docx']:
                return await self._extract_docx_content(file_path)
            elif extension in ['.txt', '.md', '.rtf']:
                return await self._extract_plain_text(file_path)
            else:
                # Fallback to plain text
                return await self._extract_plain_text(file_path)
                
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return {'text': '', 'error': str(e)}
    
    async def _extract_pdf_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PDF with layout understanding"""
        try:
            # Use PyMuPDF for better text extraction
            doc = fitz.open(str(file_path))
            
            text_parts = []
            metadata = {
                'pages': len(doc),
                'method': 'pymupdf',
                'confidence': 1.0
            }
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"Page {page_num + 1}:\n{text}")
                
                # Extract images if text is sparse
                if len(text.strip()) < 50:
                    images = page.get_images()
                    for img_index, img in enumerate(images):
                        try:
                            # Extract image and process with OCR
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            
                            if pix.n < 5:  # GRAY or RGB
                                img_data = pix.tobytes("png")
                                
                                # Process with OCR
                                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                                    tmp.write(img_data)
                                    tmp.flush()
                                    
                                    ocr_result = await self.multimodal_processor.process_image(tmp.name)
                                    if ocr_result.get('text'):
                                        text_parts.append(f"Page {page_num + 1} (OCR):\n{ocr_result['text']}")
                                    
                                    os.unlink(tmp.name)
                            
                            pix = None
                            
                        except Exception as e:
                            logger.warning(f"Failed to extract image from PDF: {e}")
            
            doc.close()
            
            return {
                'text': '\n\n'.join(text_parts),
                'metadata': metadata,
                'method': 'pymupdf_enhanced',
                'confidence': 0.95
            }
            
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return {'text': '', 'error': str(e)}
    
    async def _extract_docx_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from DOCX files"""
        try:
            doc = DocxDocument(str(file_path))
            
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(' | '.join(row_data))
                
                if table_data:
                    text_parts.append('Table:\n' + '\n'.join(table_data))
            
            return {
                'text': '\n\n'.join(text_parts),
                'method': 'python-docx',
                'confidence': 1.0
            }
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return {'text': '', 'error': str(e)}
    
    async def _extract_plain_text(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                'text': content,
                'method': 'plain_text',
                'confidence': 1.0
            }
            
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read()
                
                return {
                    'text': content,
                    'method': 'plain_text_latin1',
                    'confidence': 0.9
                }
            except Exception as e:
                logger.error(f"Plain text extraction failed: {e}")
                return {'text': '', 'error': str(e)}
    
    async def _extract_spreadsheet_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from spreadsheet files"""
        try:
            # Read spreadsheet
            if file_path.suffix.lower() == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
            
            text_parts = []
            
            if isinstance(df, dict):  # Multiple sheets
                for sheet_name, sheet_df in df.items():
                    sheet_text = self._dataframe_to_text(sheet_df, sheet_name)
                    text_parts.append(sheet_text)
            else:  # Single sheet or CSV
                sheet_text = self._dataframe_to_text(df, file_path.stem)
                text_parts.append(sheet_text)
            
            return {
                'text': '\n\n'.join(text_parts),
                'method': 'pandas',
                'confidence': 1.0
            }
            
        except Exception as e:
            logger.error(f"Spreadsheet extraction failed: {e}")
            return {'text': '', 'error': str(e)}
    
    def _dataframe_to_text(self, df: pd.DataFrame, sheet_name: str = "Sheet") -> str:
        """Convert DataFrame to structured text"""
        try:
            text_parts = [f"Sheet: {sheet_name}"]
            
            # Add column headers
            headers = ' | '.join(str(col) for col in df.columns)
            text_parts.append(f"Columns: {headers}")
            
            # Add data rows (limit to prevent overwhelming)
            max_rows = 100
            for i, row in df.head(max_rows).iterrows():
                row_text = ' | '.join(str(val) for val in row.values)
                text_parts.append(row_text)
            
            if len(df) > max_rows:
                text_parts.append(f"... ({len(df) - max_rows} more rows)")
            
            # Add summary statistics for numeric columns
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text_parts.append("\nSummary Statistics:")
                for col in numeric_cols:
                    stats = df[col].describe()
                    text_parts.append(f"{col}: mean={stats['mean']:.2f}, std={stats['std']:.2f}")
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"DataFrame to text conversion failed: {e}")
            return f"Sheet: {sheet_name}\nError processing data: {str(e)}"
    
    async def _extract_presentation_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from presentation files"""
        try:
            prs = Presentation(str(file_path))
            
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"Slide {slide_num}:"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += f"\n{shape.text}"
                
                # Extract text from tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = ' | '.join(cell.text for cell in row.cells)
                            slide_text += f"\n{row_text}"
                
                if slide_text != f"Slide {slide_num}:":
                    text_parts.append(slide_text)
            
            return {
                'text': '\n\n'.join(text_parts),
                'method': 'python-pptx',
                'confidence': 1.0
            }
            
        except Exception as e:
            logger.error(f"Presentation extraction failed: {e}")
            return {'text': '', 'error': str(e)}
    
    def _calculate_confidence_scores(self, knowledge_results: Dict[str, Any]) -> Dict[str, float]:
        """Calculate overall confidence scores for extraction results"""
        scores = {}
        
        try:
            # Entity extraction confidence
            entities = knowledge_results.get('entities', [])
            if entities:
                entity_confidences = [e.confidence for e in entities if hasattr(e, 'confidence')]
                scores['entities'] = sum(entity_confidences) / len(entity_confidences) if entity_confidences else 0.0
            else:
                scores['entities'] = 0.0
            
            # Process extraction confidence
            processes = knowledge_results.get('processes', [])
            if processes:
                # Implement process confidence calculation
                scores['processes'] = 0.8  # Placeholder
            else:
                scores['processes'] = 0.0
            
            # Relationship extraction confidence
            relationships = knowledge_results.get('relationships', [])
            if relationships:
                rel_confidences = [r.confidence for r in relationships if hasattr(r, 'confidence')]
                scores['relationships'] = sum(rel_confidences) / len(rel_confidences) if rel_confidences else 0.0
            else:
                scores['relationships'] = 0.0
            
            # Overall confidence
            all_scores = [score for score in scores.values() if score > 0]
            scores['overall'] = sum(all_scores) / len(all_scores) if all_scores else 0.0
            
        except Exception as e:
            logger.error(f"Error calculating confidence scores: {e}")
            scores = {'entities': 0.0, 'processes': 0.0, 'relationships': 0.0, 'overall': 0.0}
        
        return scores
    
    async def process_batch_documents(self, file_paths: List[str], 
                                    company_context: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """Process multiple documents in optimized batches"""
        try:
            # Create processing tasks
            tasks = []
            for i, file_path in enumerate(file_paths):
                task = self.process_document(file_path, i, company_context)
                tasks.append(task)
            
            # Process in batches
            results = await self.batch_processor.process_batch(
                tasks, 
                processor_func=self._execute_task_batch
            )
            
            return results
            
        except Exception as e:
            logger.error(f"Batch processing failed: {e}")
            return [{'error': str(e)} for _ in file_paths]
    
    async def _execute_task_batch(self, tasks: List, **kwargs) -> List[Dict[str, Any]]:
        """Execute a batch of processing tasks"""
        results = []
        for task in tasks:
            try:
                result = await task
                results.append(result)
            except Exception as e:
                logger.error(f"Task execution failed: {e}")
                results.append({'error': str(e)})
        
        return results
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get list of supported file formats"""
        return self.supported_formats
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics and model information"""
        return {
            'supported_formats': self.supported_formats,
            'model_stats': self.knowledge_engine.model_manager.get_model_stats(),
            'multimodal_capabilities': {
                'ocr_available': self.multimodal_processor.tesseract_available or self.multimodal_processor.easyocr_available,
                'audio_available': self.multimodal_processor.whisper_available,
                'vision_available': self.multimodal_processor.vision_available,
                'diarization_available': self.multimodal_processor.diarization_available
            }
        }