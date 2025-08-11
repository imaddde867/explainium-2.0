"""
EXPLAINIUM - Consolidated Document Processor

A clean, professional implementation that consolidates all document processing
functionality with support for multiple formats and AI-powered knowledge extraction.
"""

import os
import io
import json
import logging
import re
from typing import Dict, Any, List, Optional, Tuple, Union
from pathlib import Path
import tempfile
from datetime import datetime

# Core processing libraries
import requests
import cv2
import numpy as np
from PIL import Image, ImageEnhance
import pytesseract
import pandas as pd
from pptx import Presentation
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument

# AI and NLP libraries
import spacy
import torch

# Internal imports
from src.logging_config import get_logger, log_processing_step
from src.core.config import config as config_manager
from src.exceptions import ProcessingError, AIError
from src.ai.advanced_knowledge_engine import AdvancedKnowledgeEngine
from src.ai.llm_processing_engine import LLMProcessingEngine

logger = get_logger(__name__)


class DocumentProcessor:
    """
    Consolidated document processor supporting multiple formats with AI-powered knowledge extraction.
    
    Supported formats:
    - Text: PDF, DOC, DOCX, TXT, RTF
    - Images: JPG, PNG, GIF, BMP, TIFF (with OCR)
    - Spreadsheets: XLS, XLSX, CSV
    - Presentations: PPT, PPTX
    - Audio: MP3, WAV (with transcription)
    - Video: MP4, AVI (with audio extraction)
    """
    
    def __init__(self, db_session=None):
        self.tika_url = config_manager.get_tika_url()
        self.advanced_engine = AdvancedKnowledgeEngine(config_manager.ai, db_session)
        self.db_session = db_session
        self.supported_formats = {
            'text': ['.pdf', '.doc', '.docx', '.txt', '.rtf'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'],
            'spreadsheet': ['.xls', '.xlsx', '.csv'],
            'presentation': ['.ppt', '.pptx'],
            'audio': ['.mp3', '.wav', '.flac', '.aac'],
            'video': ['.mp4', '.avi', '.mov', '.mkv']
        }
        
        # Initialize OCR
        self._init_ocr()
        
        # Initialize audio processing
        self._init_audio_processing()
        
        # Initialize advanced knowledge engine
        self._init_knowledge_engine()
        
        # Initialize LLM-first processing engine
        self._init_llm_processing_engine()
        
        # Initialize AI engines  
        self.knowledge_engine_available = True
    
    def _init_ocr(self):
        """Initialize OCR capabilities"""
        try:
            # Test pytesseract
            pytesseract.get_tesseract_version()
            self.ocr_available = True
            logger.info("OCR initialized successfully")
        except Exception as e:
            logger.warning(f"OCR initialization failed: {e}")
            self.ocr_available = False
    
    def _init_audio_processing(self):
        """Initialize audio processing capabilities"""
        try:
            # Check if whisper is available
            import whisper
            # Prefer configured model name when available
            model_name = getattr(getattr(config_manager, 'ai', object()), 'whisper_model', 'base') or 'base'
            self.whisper_model = whisper.load_model(model_name)
            self.audio_processing_available = True
            logger.info("Audio processing initialized successfully")
        except Exception as e:
            logger.warning(f"Audio processing initialization failed: {e}")
            self.audio_processing_available = False
            self.whisper_model = None
    
    def _init_knowledge_engine(self):
        """Initialize the advanced knowledge engine"""
        try:
            import asyncio
            # If we're already inside an event loop (e.g., FastAPI with uvicorn reload), schedule task
            try:
                loop = asyncio.get_running_loop()
                loop.create_task(self.advanced_engine.initialize())
                logger.info("Scheduled async initialization of advanced knowledge engine")
            except RuntimeError:
                # No running loop; safe to run synchronously
                asyncio.run(self.advanced_engine.initialize())
            self.knowledge_engine_available = True  # will be true once initialize completes
        except Exception as e:
            logger.warning(f"Advanced knowledge engine initialization failed: {e}")
            self.knowledge_engine_available = False
    
    def _init_llm_processing_engine(self):
        """Initialize LLM-first processing engine"""
        try:
            import asyncio
            logger.info("Initializing LLM-first processing engine")
            self.llm_processing_engine = LLMProcessingEngine()
            
            # Schedule async initialization
            if hasattr(asyncio, 'get_running_loop'):
                try:
                    loop = asyncio.get_running_loop()
                    loop.create_task(self.llm_processing_engine.initialize())
                except RuntimeError:
                    # No running loop; safe to run synchronously
                    asyncio.run(self.llm_processing_engine.initialize())
            else:
                # No running loop; safe to run synchronously
                asyncio.run(self.llm_processing_engine.initialize())
            
            self.llm_engine_available = True
            logger.info("✅ LLM-first processing engine initialized successfully")
        except Exception as e:
            logger.error(f"❌ LLM processing engine initialization failed: {e}")
            self.llm_engine_available = False
            self.llm_processing_engine = None
    
    def _convert_llm_result_to_knowledge(self, processing_result, filename: str) -> Dict[str, Any]:
        """Convert LLM processing result to knowledge format"""
        try:
            from src.ai.knowledge_categorization_engine import KnowledgeCategory, EntityType
            
            # Convert entities to the expected format
            extracted_entities = []
            for entity in processing_result.entities:
                
                # Map to knowledge categorization format
                knowledge_entity = {
                    'entity_type': EntityType.SPECIFICATION,  # Default, will be refined
                    'key_identifier': entity.content[:50] + "..." if len(entity.content) > 50 else entity.content,
                    'core_content': entity.content,
                    'context_tags': entity.metadata.get('context_tags', []),
                    'priority_level': entity.metadata.get('priority', 'medium'),
                    'category': entity.category,
                    'confidence_score': entity.confidence,
                    'source_section': entity.context[:100] if entity.context else "",
                    'extraction_method': entity.metadata.get('extraction_method', 'llm_primary'),
                    'relationships': entity.relationships,
                    'structured_data': entity.metadata,
                    'completeness_score': entity.metadata.get('completeness', 0.8),
                    'clarity_score': entity.metadata.get('clarity', 0.8),
                    'actionability_score': entity.metadata.get('actionability', 0.7)
                }
                
                extracted_entities.append(knowledge_entity)
            
            # Build comprehensive knowledge structure
            knowledge = {
                'extracted_entities': extracted_entities,
                'processing_metadata': {
                    'method': 'llm_first_processing',
                    'confidence_score': processing_result.confidence_score,
                    'quality_metrics': processing_result.quality_metrics,
                    'processing_time': processing_result.processing_time,
                    'validation_passed': processing_result.validation_passed,
                    'llm_enhanced': processing_result.llm_enhanced,
                    'entity_count': len(extracted_entities),
                    'filename': filename,
                    'timestamp': datetime.now().isoformat()
                },
                'intelligence_framework': {
                    'document_intelligence': {
                        'document_type': processing_result.metadata.get('document_type', 'unknown'),
                        'complexity_score': processing_result.metadata.get('complexity_score', 0.5),
                        'confidence': processing_result.confidence_score
                    },
                    'processing_rules_applied': processing_result.metadata.get('rules_applied', []),
                    'quality_validation': {
                        'validation_passed': processing_result.validation_passed,
                        'quality_score': processing_result.quality_metrics.get('overall_confidence', 0.0),
                        'business_value': processing_result.quality_metrics.get('business_value', 0.0)
                    }
                },
                'quality_metrics': processing_result.quality_metrics
            }
            
            logger.info(f"✅ Converted LLM result: {len(extracted_entities)} entities")
            return knowledge
            
        except Exception as e:
            logger.error(f"❌ Failed to convert LLM result: {e}")
            return {}
    
    async def process_document_with_context(self, document: Dict[str, Any], company_context: Dict[str, Any]) -> Dict[str, Any]:
        """
        Process document with contextual understanding using company knowledge base.
        
        This method implements intelligent document processing that:
        1. Identifies document type and purpose
        2. Extracts structured and unstructured data
        3. Applies domain-specific extraction templates
        4. Cross-references with company knowledge base
        5. Generates semantic embeddings for similarity search
        """
        try:
            if not self.advanced_engine:
                raise AIError("Advanced knowledge engine not initialized")
            
            # Extract document metadata and content
            doc_type = document.get('type', 'unknown')
            content = document.get('content', '')
            metadata = document.get('metadata', {})
            
            # Step 1: Identify document type and purpose
            document_purpose = await self._identify_document_purpose(content, doc_type, company_context)
            
            # Step 2: Extract structured and unstructured data
            structured_data = await self._extract_structured_data(content, doc_type, document_purpose)
            unstructured_data = await self._extract_unstructured_data(content, doc_type, document_purpose)
            
            # Step 3: Apply domain-specific extraction templates
            domain_data = await self._apply_domain_templates(content, doc_type, company_context)
            
            # Step 4: Cross-reference with company knowledge base
            cross_references = await self._cross_reference_knowledge(content, company_context)
            
            # Step 5: Generate semantic embeddings
            embeddings = await self._generate_semantic_embeddings(content, structured_data)
            
            # Combine all extracted information
            enhanced_document = {
                'original_document': document,
                'document_purpose': document_purpose,
                'structured_data': structured_data,
                'unstructured_data': unstructured_data,
                'domain_data': domain_data,
                'cross_references': cross_references,
                'embeddings': embeddings,
                'processing_timestamp': datetime.now().isoformat(),
                'processing_method': 'intelligent_contextual'
            }
            
            # Extract deep knowledge using the advanced engine
            knowledge_extraction = await self.advanced_engine.extract_deep_knowledge(enhanced_document)
            enhanced_document['knowledge_extraction'] = knowledge_extraction
            
            logger.info(f"Successfully processed document with context: {document.get('id', 'unknown')}")
            return enhanced_document
            
        except Exception as e:
            logger.error(f"Failed to process document with context: {e}")
            raise ProcessingError(f"Contextual processing failed: {str(e)}")
    
    async def extract_tacit_knowledge(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Extract tacit knowledge and patterns across multiple documents.
        
        This method implements pattern recognition that identifies:
        - Recurring themes and patterns
        - Implicit workflows from email chains
        - Organizational structures from documents
        - Policy changes over time
        - Informal communication networks
        """
        try:
            if not self.advanced_engine:
                raise AIError("Advanced knowledge engine not initialized")
            
            # Step 1: Analyze document patterns
            document_patterns = await self._analyze_document_patterns(documents)
            
            # Step 2: Extract implicit workflows
            implicit_workflows = await self._extract_implicit_workflows(documents)
            
            # Step 3: Identify organizational structures
            org_structures = await self._identify_organizational_structures(documents)
            
            # Step 4: Detect policy changes
            policy_changes = await self._detect_policy_changes(documents)
            
            # Step 5: Map communication networks
            communication_networks = await self._map_communication_networks(documents)
            
            # Combine all tacit knowledge insights
            tacit_knowledge = {
                'document_patterns': document_patterns,
                'implicit_workflows': implicit_workflows,
                'organizational_structures': org_structures,
                'policy_changes': policy_changes,
                'communication_networks': communication_networks,
                'extraction_timestamp': datetime.now().isoformat(),
                'total_documents_analyzed': len(documents)
            }
            
            # Use advanced engine to build knowledge graph from tacit knowledge
            knowledge_graph = await self.advanced_engine.build_knowledge_graph(tacit_knowledge)
            tacit_knowledge['knowledge_graph'] = knowledge_graph
            
            logger.info(f"Successfully extracted tacit knowledge from {len(documents)} documents")
            return tacit_knowledge
                
        except Exception as e:
            logger.error(f"Failed to extract tacit knowledge: {e}")
            raise ProcessingError(f"Tacit knowledge extraction failed: {str(e)}")
    
    def process_document(self, file_path: str, document_id: int) -> Dict[str, Any]:
        """
        Main method to process any supported document format
        
        Args:
            file_path: Path to the document file
            document_id: Database ID of the document
            
        Returns:
            Dictionary containing extracted content and knowledge
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise ProcessingError(f"File not found: {file_path}")
        
        file_extension = file_path.suffix.lower()
        file_type = self._get_file_type(file_extension)
        
        logger.info(f"Processing {file_type} document: {file_path.name}")
        
        try:
            # Extract content based on file type
            if file_type == 'text':
                content = self._process_text_document(file_path)
            elif file_type == 'image':
                content = self._process_image_document(file_path)
            elif file_type == 'spreadsheet':
                content = self._process_spreadsheet_document(file_path)
            elif file_type == 'presentation':
                content = self._process_presentation_document(file_path)
            elif file_type == 'audio':
                content = self._process_audio_document(file_path)
            elif file_type == 'video':
                content = self._process_video_document(file_path)
            else:
                raise ProcessingError(f"Unsupported file type: {file_extension}")
            
            # PRIMARY: LLM-First Processing Engine (Superior Results)
            knowledge = {}
            
            # RULE 1: LLM-First Processing (Primary Method - Guaranteed Best Results)
            if hasattr(self, 'llm_engine_available') and self.llm_engine_available and self.llm_processing_engine and self.llm_processing_engine.initialized:
                try:
                    logger.info("🧠 Using LLM-First Processing Engine (Primary Method)")
                    import asyncio
                    loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(loop)
                    
                    # Process with LLM-first engine
                    processing_result = loop.run_until_complete(
                        self.llm_processing_engine.process_document(
                            content=content.get('text', ''),
                            document_type=file_type,
                            metadata={
                                'filename': file_path.name,
                                'file_type': file_type,
                                'document_id': document_id,
                                'content_type': content.get('content_type', 'text'),
                                'processing_timestamp': datetime.now().isoformat()
                            }
                        )
                    )
                    loop.close()
                    
                    # Convert LLM processing result to knowledge format
                    knowledge = self._convert_llm_result_to_knowledge(processing_result, file_path.name)
                    
                    # Set processing method metadata
                    if 'processing_metadata' not in knowledge:
                        knowledge['processing_metadata'] = {}
                    knowledge['processing_metadata']['processing_method'] = 'llm_primary'
                    knowledge['extraction_methods'] = ['llm_primary']
                    
                    logger.info(f"✅ LLM-First processing successful: {len(processing_result.entities)} entities, "
                              f"{processing_result.confidence_score:.2f} confidence")
                    
                except Exception as e:
                    logger.warning(f"⚠️ LLM-First processing failed, cascading to fallback: {e}")
                    knowledge = {}
            
            # FALLBACK 1: Advanced Knowledge Engine (High Quality Fallback)
            if not knowledge and self.knowledge_engine_available:
                try:
                    logger.info("🔧 Using Advanced Knowledge Engine (Fallback Method)")
                    import asyncio
                    loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(loop)
                    
                    document_data = {
                        'id': document_id,
                        'content': content.get('text', ''),
                        'filename': file_path.name,
                        'metadata': {
                            'filename': file_path.name,
                            'file_type': file_type,
                            'document_id': document_id,
                            'content_type': content.get('content_type', 'text'),
                            'processing_timestamp': datetime.now().isoformat(),
                            'fallback_reason': 'llm_engine_unavailable'
                        },
                        'sections': content.get('sections', []),
                        'extraction_methods': content.get('extraction_methods', [])
                    }
                    
                    knowledge_results = loop.run_until_complete(
                        self.advanced_engine.extract_intelligent_knowledge(document_data)
                    )
                    knowledge = knowledge_results
                    loop.close()
                    logger.info("✅ Advanced knowledge extraction successful (fallback)")
                    
                except Exception as e:
                    logger.error(f"❌ Advanced knowledge extraction failed, using legacy: {e}")
                    knowledge = {}
            
            # FALLBACK 2: Legacy Processing (Basic)
            if not knowledge:
                logger.warning("⚠️ Using legacy processing (basic patterns)")
                knowledge = self._extract_knowledge_legacy(content.get('text', ''), file_path.name)
            
            result = {
                'document_id': document_id,
                'filename': file_path.name,
                'file_type': file_type,
                'processed_at': datetime.now().isoformat(),
                'content': content,
                'knowledge': knowledge,
                'processing_status': 'completed'
            }
            
            logger.info(f"Successfully processed document: {file_path.name}")
            return result
            
        except Exception as e:
            logger.error(f"Error processing document {file_path.name}: {e}")
            raise ProcessingError(f"Failed to process document: {str(e)}")
    
    def _get_file_type(self, extension: str) -> str:
        """Determine file type from extension"""
        for file_type, extensions in self.supported_formats.items():
            if extension in extensions:
                return file_type
        return 'unknown'
    
    def _process_text_document(self, file_path: Path) -> Dict[str, Any]:
        """Process text-based documents (PDF, DOC, DOCX, TXT)"""
        extension = file_path.suffix.lower()
        
        if extension == '.pdf':
            return self._extract_pdf_content(file_path)
        elif extension in ['.doc', '.docx']:
            return self._extract_word_content(file_path)
        elif extension == '.txt':
            return self._extract_text_content(file_path)
        else:
            # Try Tika for other formats
            return self._extract_with_tika(file_path)
    
    def _extract_pdf_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PDF files"""
        text_content = ""
        metadata = {}
        
        try:
            # Try PyMuPDF first (better for complex PDFs)
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text_content += page.get_text()
            
            metadata = doc.metadata
            doc.close()
            
        except Exception as e:
            logger.warning(f"PyMuPDF failed, trying PyPDF2: {e}")
            try:
                # Fallback to PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text_content += page.extract_text()
                    
                    if pdf_reader.metadata:
                        metadata = dict(pdf_reader.metadata)
            except Exception as e2:
                logger.error(f"Both PDF extraction methods failed: {e2}")
                raise ProcessingError(f"Failed to extract PDF content: {e2}")
        
        # Extract sections for intelligent processing
        sections = self._extract_document_sections(text_content)
        
        return {
            'text': text_content.strip(),
            'metadata': metadata,
            'page_count': len(text_content.split('\f')) if '\f' in text_content else 1,
            'sections': sections
        }
    
    def _extract_word_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Word documents"""
        try:
            doc = DocxDocument(file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract tables
            table_content = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                table_content.append(table_data)
            
            return {
                'text': text_content.strip(),
                'tables': table_content,
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables)
            }
            
        except Exception as e:
            logger.error(f"Failed to extract Word content: {e}")
            raise ProcessingError(f"Failed to extract Word content: {e}")
    
    def _extract_text_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            return {
                'text': content,
                'line_count': len(content.split('\n')),
                'character_count': len(content)
            }
            
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    return {
                        'text': content,
                        'encoding': encoding,
                        'line_count': len(content.split('\n')),
                        'character_count': len(content)
                    }
                except UnicodeDecodeError:
                    continue
            
            raise ProcessingError("Unable to decode text file with any supported encoding")
    
    def _process_image_document(self, file_path: Path) -> Dict[str, Any]:
        """Process image documents with advanced OCR and computer vision analysis"""
        try:
            # Load and preprocess image
            image = cv2.imread(str(file_path))
            if image is None:
                raise ProcessingError("Unable to load image")
            
            # Get image metadata
            height, width, channels = image.shape
            
            # Multi-layered text extraction
            extracted_texts = []
            extraction_methods = []
            
            # Method 1: Enhanced OCR with preprocessing
            if self.ocr_available:
                ocr_text = self._enhanced_ocr_extraction(image)
                if ocr_text.strip():
                    extracted_texts.append(ocr_text)
                    extraction_methods.append("enhanced_ocr")
            
            # Method 2: Computer vision analysis
            cv_analysis = self._analyze_image_content(image, file_path)
            if cv_analysis.get('description'):
                extracted_texts.append(cv_analysis['description'])
                extraction_methods.append("computer_vision")
            
            # Method 3: Document structure analysis
            if self.ocr_available:
                structure_analysis = self._analyze_document_structure_in_image(image)
                if structure_analysis.get('structural_text'):
                    extracted_texts.append(structure_analysis['structural_text'])
                    extraction_methods.append("structure_analysis")
            
            # Method 4: Diagram and chart analysis
            diagram_analysis = self._analyze_diagrams_and_charts(image)
            if diagram_analysis.get('diagram_text'):
                extracted_texts.append(diagram_analysis['diagram_text'])
                extraction_methods.append("diagram_analysis")
            
            # Combine all extracted content
            combined_text = self._combine_image_extractions(extracted_texts, extraction_methods)
            
            # Extract sections for intelligent processing
            sections = self._extract_document_sections(combined_text) if combined_text else []
            
            result = {
                'text': combined_text,
                'image_width': width,
                'image_height': height,
                'image_channels': channels,
                'extraction_methods': extraction_methods,
                'sections': sections,
                'computer_vision_analysis': cv_analysis,
                'structure_analysis': structure_analysis if 'structure_analysis' in locals() else {},
                'diagram_analysis': diagram_analysis,
                'ocr_confidence': self._calculate_ocr_confidence(combined_text) if combined_text else 0.0
            }
            
            # Add intelligent metadata
            result['content_type_detected'] = self._detect_image_content_type(image, combined_text)
            result['visual_elements'] = self._identify_visual_elements(image)
            
            return result
            
        except Exception as e:
            logger.error(f"Failed to process image: {e}")
            # Fallback to basic processing
            try:
                return self._basic_image_processing_fallback(file_path)
            except:
                raise ProcessingError(f"Failed to process image: {e}")
    
    def _process_spreadsheet_document(self, file_path: Path) -> Dict[str, Any]:
        """Process spreadsheet documents"""
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
            
            if isinstance(df, dict):  # Multiple sheets
                text_content = ""
                sheet_info = {}
                for sheet_name, sheet_df in df.items():
                    sheet_text = self._dataframe_to_text(sheet_df)
                    text_content += f"\n\n=== Sheet: {sheet_name} ===\n{sheet_text}"
                    sheet_info[sheet_name] = {
                        'rows': len(sheet_df),
                        'columns': len(sheet_df.columns)
                    }
                
                return {
                    'text': text_content.strip(),
                    'sheets': sheet_info,
                    'total_sheets': len(df)
                }
            else:  # Single sheet
                text_content = self._dataframe_to_text(df)
                return {
                    'text': text_content,
                    'rows': len(df),
                    'columns': len(df.columns)
                }
                
        except Exception as e:
            logger.error(f"Failed to process spreadsheet: {e}")
            raise ProcessingError(f"Failed to process spreadsheet: {e}")
    
    def _process_presentation_document(self, file_path: Path) -> Dict[str, Any]:
        """Process presentation documents"""
        try:
            prs = Presentation(file_path)
            text_content = ""
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                slide_text = f"\n\n=== Slide {slide_count} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                text_content += slide_text
            
            return {
                'text': text_content.strip(),
                'slide_count': slide_count
            }
            
        except Exception as e:
            logger.error(f"Failed to process presentation: {e}")
            raise ProcessingError(f"Failed to process presentation: {e}")
    
    def _process_audio_document(self, file_path: Path) -> Dict[str, Any]:
        """Process audio documents with transcription"""
        if not self.audio_processing_available:
            raise ProcessingError("Audio processing not available")
        
        try:
            # Transcribe audio using Whisper
            result = self.whisper_model.transcribe(str(file_path))
            
            return {
                'text': result['text'],
                'language': result.get('language', 'unknown'),
                'segments': result.get('segments', [])
            }
            
        except Exception as e:
            logger.error(f"Failed to process audio: {e}")
            raise ProcessingError(f"Failed to process audio: {e}")
    
    def _process_video_document(self, file_path: Path) -> Dict[str, Any]:
        """Enhanced video processing with advanced frame analysis and audio transcription.

        Multi-layered strategy:
        1) Audio transcription with Whisper for spoken content
        2) Advanced frame sampling with intelligent scene detection
        3) On-screen text extraction with enhanced OCR
        4) Visual content analysis for procedural videos
        5) Scene segmentation and content categorization
        """
        # Initialize extraction results
        extracted_content = {
            'audio_transcript': '',
            'visual_text': '',
            'scene_analysis': '',
            'procedural_content': '',
            'metadata': {}
        }
        
        language_detected = "unknown"
        segments = []
        sources = []
        processing_metadata = {}

        # Get video metadata
        try:
            video_metadata = self._extract_video_metadata(file_path)
            processing_metadata.update(video_metadata)
        except Exception as e:
            logger.warning(f"Failed to extract video metadata: {e}")

        # Method 1: Enhanced Audio Transcription
        if self.audio_processing_available and self.whisper_model is not None:
            try:
                audio_result = self._enhanced_audio_transcription(file_path)
                if audio_result.get('transcript'):
                    extracted_content['audio_transcript'] = audio_result['transcript']
                    language_detected = audio_result.get('language', 'unknown')
                    segments = audio_result.get('segments', [])
                    sources.append('enhanced_audio_transcription')
                    processing_metadata['audio_processing'] = audio_result.get('metadata', {})
            except Exception as e:
                logger.warning(f"Enhanced audio transcription failed: {e}")

        # Method 2: Intelligent Frame Analysis
        try:
            frame_analysis = self._intelligent_frame_analysis(file_path)
            if frame_analysis.get('visual_text'):
                extracted_content['visual_text'] = frame_analysis['visual_text']
                sources.append('intelligent_frame_analysis')
            if frame_analysis.get('scene_analysis'):
                extracted_content['scene_analysis'] = frame_analysis['scene_analysis']
            processing_metadata['frame_analysis'] = frame_analysis.get('metadata', {})
        except Exception as e:
            logger.warning(f"Intelligent frame analysis failed: {e}")

        # Method 3: Procedural Content Detection
        try:
            procedural_analysis = self._detect_procedural_content(file_path, extracted_content)
            if procedural_analysis.get('procedural_text'):
                extracted_content['procedural_content'] = procedural_analysis['procedural_text']
                sources.append('procedural_analysis')
            processing_metadata['procedural_analysis'] = procedural_analysis.get('metadata', {})
        except Exception as e:
            logger.warning(f"Procedural content detection failed: {e}")

        # Method 4: Scene Segmentation and Categorization
        try:
            scene_segmentation = self._analyze_video_scenes(file_path)
            if scene_segmentation.get('scene_descriptions'):
                extracted_content['scene_analysis'] = scene_segmentation['scene_descriptions']
                sources.append('scene_segmentation')
            processing_metadata['scene_segmentation'] = scene_segmentation.get('metadata', {})
        except Exception as e:
            logger.warning(f"Scene segmentation failed: {e}")

        # Combine all extracted content intelligently
        combined_text = self._combine_video_extractions(extracted_content, sources)
        
        # Extract sections for intelligent processing
        sections = self._extract_document_sections(combined_text) if combined_text else []

        # Fallback to basic processing if nothing extracted
        if not combined_text.strip():
            try:
                basic_result = self._basic_video_processing_fallback(file_path)
                combined_text = basic_result.get('text', '')
                sources.append('basic_fallback')
                processing_metadata['fallback_used'] = True
            except Exception as e:
                logger.warning(f"Even basic video processing failed: {e}")

        # Compile comprehensive results
        result = {
            'text': combined_text,
            'language': language_detected,
            'segments': segments,
            'sources': sources,
            'sections': sections,
            'extraction_methods': sources,
            'video_metadata': processing_metadata,
            'content_breakdown': extracted_content,
            'video_analysis': {
                'duration': processing_metadata.get('duration', 0),
                'fps': processing_metadata.get('fps', 0),
                'frames_analyzed': processing_metadata.get('frames_analyzed', 0),
                'audio_quality': processing_metadata.get('audio_quality', 'unknown'),
                'visual_quality': processing_metadata.get('visual_quality', 'unknown')
            }
        }

        # Add intelligent metadata
        result['content_type_detected'] = self._detect_video_content_type(combined_text, processing_metadata)
        result['knowledge_extraction_confidence'] = self._calculate_video_extraction_confidence(extracted_content, sources)

        return result

    def _extract_text_from_video_frames(self, file_path: Path, interval_seconds: int = 5, max_frames: int = 20) -> Dict[str, Any]:
        """Extract visible text from sampled video frames using OCR.

        Samples one frame approximately every `interval_seconds`, up to `max_frames` frames.
        Applies simple preprocessing to improve OCR accuracy.
        """
        capture = cv2.VideoCapture(str(file_path))
        if not capture.isOpened():
            raise ProcessingError("Unable to open video file for frame OCR")

        try:
            fps = capture.get(cv2.CAP_PROP_FPS) or 0.0
            total_frames = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
            frames_between_samples = max(1, int(round((fps or 1.0) * max(1, interval_seconds))))

            collected_lines: List[str] = []
            frames_analyzed = 0
            frame_index = 0

            while frame_index < total_frames and frames_analyzed < max_frames:
                capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
                success, frame = capture.read()
                if not success or frame is None:
                    break

                gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                # Light denoising and thresholding to enhance text regions
                denoised = cv2.fastNlMeansDenoising(gray)
                _, thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

                text = pytesseract.image_to_string(thresh) if self.ocr_available else ''
                if text and text.strip():
                    # Keep unique non-empty lines to reduce noise
                    for line in [l.strip() for l in text.splitlines() if l.strip()]:
                        if line not in collected_lines:
                            collected_lines.append(line)

                frames_analyzed += 1
                frame_index += frames_between_samples

            ocr_text = "\n".join(collected_lines)
            return {'text': ocr_text, 'frames_analyzed': frames_analyzed}
        finally:
            capture.release()
    
    def _extract_with_tika(self, file_path: Path) -> Dict[str, Any]:
        """Extract content using Apache Tika server"""
        if not self.tika_url:
            raise ProcessingError("Tika server not configured")
        
        try:
            with open(file_path, 'rb') as file:
                response = requests.put(
                    f"{self.tika_url}/tika",
                    data=file,
                    headers={'Accept': 'text/plain'},
                    timeout=30
                )
                response.raise_for_status()
                
                return {
                    'text': response.text.strip(),
                    'extraction_method': 'tika'
                }
                
        except Exception as e:
            logger.error(f"Tika extraction failed: {e}")
            raise ProcessingError(f"Tika extraction failed: {e}")
    
    def _dataframe_to_text(self, df: pd.DataFrame) -> str:
        """Convert DataFrame to readable text"""
        text_parts = []
        
        # Add column headers
        text_parts.append("Columns: " + ", ".join(df.columns.astype(str)))
        
        # Add data rows (limit to prevent excessive text)
        max_rows = 100
        for idx, row in df.head(max_rows).iterrows():
            row_text = " | ".join([f"{col}: {val}" for col, val in row.items()])
            text_parts.append(f"Row {idx + 1}: {row_text}")
        
        if len(df) > max_rows:
            text_parts.append(f"... and {len(df) - max_rows} more rows")
        
        return "\n".join(text_parts)
    
    def _calculate_ocr_confidence(self, text: str) -> float:
        """Calculate OCR confidence based on text characteristics"""
        if not text.strip():
            return 0.0
        
        # Simple confidence calculation based on text characteristics
        confidence = 1.0
        
        # Penalize for common OCR errors
        if any(char in text for char in '|[]{}()'):
            confidence -= 0.1
        
        # Penalize for excessive numbers (common OCR artifact)
        if sum(c.isdigit() for c in text) / len(text) > 0.3:
            confidence -= 0.2
        
        # Penalize for very short words (OCR artifacts)
        words = text.split()
        if words and any(len(word) == 1 for word in words):
            confidence -= 0.1
        
        return max(0.0, confidence)
    
    # Helper methods for intelligent document processing
    
    async def _identify_document_purpose(self, content: str, doc_type: str, company_context: Dict[str, Any]) -> Dict[str, Any]:
        """Identify the purpose and intent of a document"""
        try:
            # Use LLM to analyze document purpose
            purpose_prompt = f"""
            Analyze this {doc_type} document and identify its purpose:
            
            Content: {content[:1000]}...
            
            Company Context: {company_context.get('industry', 'Unknown')} industry
            
            Identify:
            1. Primary purpose (e.g., policy, procedure, report, communication)
            2. Target audience
            3. Business function
            4. Urgency level
            5. Compliance requirements
            """
            
            purpose_analysis = await self.advanced_engine.llm.analyze(purpose_prompt)
            
            return {
                'primary_purpose': purpose_analysis.get('primary_purpose', 'unknown'),
                'target_audience': purpose_analysis.get('target_audience', 'unknown'),
                'business_function': purpose_analysis.get('business_function', 'unknown'),
                'urgency_level': purpose_analysis.get('urgency_level', 'low'),
                'compliance_requirements': purpose_analysis.get('compliance_requirements', [])
            }
            
        except Exception as e:
            logger.warning(f"Failed to identify document purpose: {e}")
            return {'primary_purpose': 'unknown', 'error': str(e)}
    
    async def _extract_structured_data(self, content: str, doc_type: str, purpose: Dict[str, Any]) -> Dict[str, Any]:
        """Extract structured data based on document type and purpose"""
        try:
            structured_data = {}
            
            if doc_type in ['spreadsheet', 'csv']:
                # Extract tables, formulas, and data relationships
                structured_data['tables'] = await self._extract_table_structures(content)
                structured_data['formulas'] = await self._extract_formulas(content)
                structured_data['data_relationships'] = await self._extract_data_relationships(content)
            
            elif doc_type in ['presentation', 'ppt', 'pptx']:
                # Extract slide structure, key points, and flow
                structured_data['slide_structure'] = await self._extract_slide_structure(content)
                structured_data['key_points'] = await self._extract_key_points(content)
                structured_data['presentation_flow'] = await self._extract_presentation_flow(content)
            
            elif doc_type in ['text', 'pdf', 'doc', 'docx']:
                # Extract sections, headings, and document structure
                structured_data['sections'] = await self._extract_document_sections(content)
                structured_data['headings'] = await self._extract_headings(content)
                structured_data['document_structure'] = await self._extract_document_structure(content)
            
            return structured_data
            
        except Exception as e:
            logger.warning(f"Failed to extract structured data: {e}")
            return {'error': str(e)}
    
    async def _extract_unstructured_data(self, content: str, doc_type: str, purpose: Dict[str, Any]) -> Dict[str, Any]:
        """Extract unstructured data and insights"""
        try:
            unstructured_data = {}
            
            # Extract key concepts and entities
            concepts = await self.advanced_engine.extract_concepts(content)
            unstructured_data['concepts'] = concepts
            
            # Extract sentiment and tone
            sentiment = await self.advanced_engine.analyze_sentiment(content)
            unstructured_data['sentiment'] = sentiment
            
            # Extract key phrases and topics
            topics = await self.advanced_engine.extract_topics(content)
            unstructured_data['topics'] = topics
            
            # Extract named entities
            entities = await self.advanced_engine.extract_entities(content)
            unstructured_data['entities'] = entities
            
            return unstructured_data
            
        except Exception as e:
            logger.warning(f"Failed to extract unstructured data: {e}")
            return {'error': str(e)}
    
    async def _apply_domain_templates(self, content: str, doc_type: str, company_context: Dict[str, Any]) -> Dict[str, Any]:
        """Apply domain-specific extraction templates"""
        try:
            domain_data = {}
            industry = company_context.get('industry', 'general')
            
            if industry == 'healthcare':
                domain_data.update(await self._extract_healthcare_entities(content))
            elif industry == 'finance':
                domain_data.update(await self._extract_finance_entities(content))
            elif industry == 'legal':
                domain_data.update(await self._extract_legal_entities(content))
            elif industry == 'manufacturing':
                domain_data.update(await self._extract_manufacturing_entities(content))
            else:
                domain_data.update(await self._extract_general_entities(content))
            
            return domain_data
            
        except Exception as e:
            logger.warning(f"Failed to apply domain templates: {e}")
            return {'error': str(e)}
    
    async def _cross_reference_knowledge(self, content: str, company_context: Dict[str, Any]) -> Dict[str, Any]:
        """Cross-reference content with existing company knowledge base"""
        try:
            cross_references = {}
            
            # Find similar documents
            similar_docs = await self.advanced_engine.find_similar_documents(content)
            cross_references['similar_documents'] = similar_docs
            
            # Find related processes
            related_processes = await self.advanced_engine.find_related_processes(content)
            cross_references['related_processes'] = related_processes
            
            # Find related systems
            related_systems = await self.advanced_engine.find_related_systems(content)
            cross_references['related_systems'] = related_systems
            
            # Find related requirements
            related_requirements = await self.advanced_engine.find_related_requirements(content)
            cross_references['related_requirements'] = related_requirements
            
            return cross_references
            
        except Exception as e:
            logger.warning(f"Failed to cross-reference knowledge: {e}")
            return {'error': str(e)}
    
    async def _generate_semantic_embeddings(self, content: str, structured_data: Dict[str, Any]) -> Dict[str, Any]:
        """Generate semantic embeddings for similarity search"""
        try:
            embeddings = {}
            
            # Generate content embedding
            content_embedding = await self.advanced_engine.generate_embedding(content)
            embeddings['content'] = content_embedding
            
            # Generate embeddings for structured data
            if 'key_points' in structured_data:
                key_points_embedding = await self.advanced_engine.generate_embedding(
                    ' '.join(structured_data['key_points'])
                )
                embeddings['key_points'] = key_points_embedding
            
            if 'topics' in structured_data:
                topics_embedding = await self.advanced_engine.generate_embedding(
                    ' '.join(structured_data['topics'])
                )
                embeddings['topics'] = topics_embedding
            
            return embeddings
            
        except Exception as e:
            logger.warning(f"Failed to generate embeddings: {e}")
            return {'error': str(e)}
    
    async def _analyze_document_patterns(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Analyze patterns across multiple documents"""
        try:
            patterns = {}
            
            # Analyze content patterns
            content_patterns = await self.advanced_engine.analyze_content_patterns(documents)
            patterns['content'] = content_patterns
            
            # Analyze metadata patterns
            metadata_patterns = await self.advanced_engine.analyze_metadata_patterns(documents)
            patterns['metadata'] = metadata_patterns
            
            # Analyze temporal patterns
            temporal_patterns = await self.advanced_engine.analyze_temporal_patterns(documents)
            patterns['temporal'] = temporal_patterns
            
            return patterns
            
        except Exception as e:
            logger.warning(f"Failed to analyze document patterns: {e}")
            return {'error': str(e)}
    
    async def _extract_implicit_workflows(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Extract implicit workflows from document patterns"""
        try:
            workflows = {}
            
            # Identify process flows
            process_flows = await self.advanced_engine.identify_process_flows(documents)
            workflows['process_flows'] = process_flows
            
            # Identify decision points
            decision_points = await self.advanced_engine.identify_decision_points(documents)
            workflows['decision_points'] = decision_points
            
            # Identify handoffs
            handoffs = await self.advanced_engine.identify_handoffs(documents)
            workflows['handoffs'] = handoffs
            
            return workflows
            
        except Exception as e:
            logger.warning(f"Failed to extract implicit workflows: {e}")
            return {'error': str(e)}
    
    async def _identify_organizational_structures(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Identify organizational structures from documents"""
        try:
            org_structures = {}
            
            # Identify roles and responsibilities
            roles = await self.advanced_engine.identify_roles_responsibilities(documents)
            org_structures['roles'] = roles
            
            # Identify reporting relationships
            reporting = await self.advanced_engine.identify_reporting_relationships(documents)
            org_structures['reporting'] = reporting
            
            # Identify teams and departments
            teams = await self.advanced_engine.identify_teams_departments(documents)
            org_structures['teams'] = teams
            
            return org_structures
            
        except Exception as e:
            logger.warning(f"Failed to identify organizational structures: {e}")
            return {'error': str(e)}
    
    async def _detect_policy_changes(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Detect policy changes over time"""
        try:
            policy_changes = {}
            
            # Identify policy evolution
            evolution = await self.advanced_engine.analyze_policy_evolution(documents)
            policy_changes['evolution'] = evolution
            
            # Identify version differences
            versions = await self.advanced_engine.analyze_version_differences(documents)
            policy_changes['versions'] = versions
            
            # Identify compliance changes
            compliance = await self.advanced_engine.analyze_compliance_changes(documents)
            policy_changes['compliance'] = compliance
            
            return policy_changes
            
        except Exception as e:
            logger.warning(f"Failed to detect policy changes: {e}")
            return {'error': str(e)}
    
    async def _map_communication_networks(self, documents: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Map informal communication networks"""
        try:
            networks = {}
            
            # Identify communication patterns
            patterns = await self.advanced_engine.identify_communication_patterns(documents)
            networks['patterns'] = patterns
            
            # Identify key communicators
            communicators = await self.advanced_engine.identify_key_communicators(documents)
            networks['communicators'] = communicators
            
            # Identify information flow
            flow = await self.advanced_engine.identify_information_flow(documents)
            networks['flow'] = flow
            
            return networks
            
        except Exception as e:
            logger.warning(f"Failed to map communication networks: {e}")
            return {'error': str(e)}
    
    # Placeholder methods for structured data extraction
    async def _extract_table_structures(self, content: str) -> List[Dict[str, Any]]:
        """Extract table structures from content"""
        # Placeholder implementation
        return []
    
    async def _extract_formulas(self, content: str) -> List[str]:
        """Extract formulas from content"""
        # Placeholder implementation
        return []
    
    async def _extract_data_relationships(self, content: str) -> List[Dict[str, Any]]:
        """Extract data relationships from content"""
        # Placeholder implementation
        return []
    
    async def _extract_slide_structure(self, content: str) -> Dict[str, Any]:
        """Extract slide structure from content"""
        # Placeholder implementation
        return {}
    
    async def _extract_key_points(self, content: str) -> List[str]:
        """Extract key points from content"""
        # Placeholder implementation
        return []
    
    async def _extract_presentation_flow(self, content: str) -> List[Dict[str, Any]]:
        """Extract presentation flow from content"""
        # Placeholder implementation
        return []
    
    async def _extract_document_sections(self, content: str) -> List[Dict[str, Any]]:
        """Extract document sections from content"""
        # Placeholder implementation
        return []
    
    async def _extract_headings(self, content: str) -> List[str]:
        """Extract headings from content"""
        # Placeholder implementation
        return []
    
    async def _extract_document_structure(self, content: str) -> Dict[str, Any]:
        """Extract document structure from content"""
        # Placeholder implementation
        return {}
    
    # Placeholder methods for domain-specific extraction
    async def _extract_healthcare_entities(self, content: str) -> Dict[str, Any]:
        """Extract healthcare-specific entities"""
        # Placeholder implementation
        return {}
    
    async def _extract_finance_entities(self, content: str) -> Dict[str, Any]:
        """Extract finance-specific entities"""
        # Placeholder implementation
        return {}
    
    async def _extract_legal_entities(self, content: str) -> Dict[str, Any]:
        """Extract legal-specific entities"""
        # Placeholder implementation
        return {}
    
    async def _extract_manufacturing_entities(self, content: str) -> Dict[str, Any]:
        """Extract manufacturing-specific entities"""
        # Placeholder implementation
        return {}
    
    async def _extract_general_entities(self, content: str) -> Dict[str, Any]:
        """Extract general entities"""
        # Placeholder implementation
        return {}
        
            # Note: The OCR confidence helper and its heuristic are defined above
    
    def _extract_document_sections(self, text: str) -> List[Dict[str, Any]]:
        """Extract document sections for intelligent processing"""
        sections = []
        
        # Split text into logical sections
        section_patterns = [
            r'\n#+\s*(.*?)\n',  # Markdown headers
            r'\n(\d+\.?\s+.*?)\n',  # Numbered sections
            r'\n([A-Z][A-Z\s]{3,})\n',  # ALL CAPS headers
            r'\n(SECTION.*?)\n',  # Section headers
            r'\n(CHAPTER.*?)\n',  # Chapter headers
        ]
        
        current_section = ""
        current_title = "Introduction"
        section_id = 0
        
        for line in text.split('\n'):
            line_stripped = line.strip()
            
            # Check if this line is a section header
            is_header = False
            for pattern in section_patterns:
                if re.match(pattern, f'\n{line}\n'):
                    is_header = True
                    # Save previous section
                    if current_section.strip():
                        sections.append({
                            'id': section_id,
                            'title': current_title,
                            'content': current_section.strip(),
                            'word_count': len(current_section.split()),
                            'type': self._classify_section_type(current_title, current_section)
                        })
                        section_id += 1
                    
                    # Start new section
                    current_title = line_stripped
                    current_section = ""
                    break
            
            if not is_header:
                current_section += line + '\n'
        
        # Add final section
        if current_section.strip():
            sections.append({
                'id': section_id,
                'title': current_title,
                'content': current_section.strip(),
                'word_count': len(current_section.split()),
                'type': self._classify_section_type(current_title, current_section)
            })
        
        return sections
    
    def _classify_section_type(self, title: str, content: str) -> str:
        """Classify section type for intelligent processing"""
        title_lower = title.lower()
        content_lower = content.lower()
        
        if any(keyword in title_lower for keyword in ["process", "procedure", "step"]):
            return "process"
        elif any(keyword in title_lower for keyword in ["requirement", "compliance", "standard"]):
            return "compliance"
        elif any(keyword in title_lower for keyword in ["risk", "safety", "hazard"]):
            return "risk"
        elif any(keyword in title_lower for keyword in ["role", "responsibility", "organization"]):
            return "organizational"
        elif any(keyword in title_lower for keyword in ["definition", "glossary", "term"]):
            return "definition"
        elif any(keyword in content_lower for keyword in ["must", "shall", "mandatory", "required"]):
            return "requirement"
        elif any(keyword in content_lower for keyword in ["step", "procedure", "process"]):
            return "process"
        else:
            return "general"
    
    # Enhanced Image Processing Methods
    
    def _enhanced_ocr_extraction(self, image: np.ndarray) -> str:
        """Enhanced OCR with multiple preprocessing techniques"""
        try:
            if not self.ocr_available:
                return ""
            
            # Convert to grayscale
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            
            # Multiple preprocessing approaches
            text_results = []
            
            # Approach 1: Standard preprocessing
            denoised = cv2.fastNlMeansDenoising(gray)
            text1 = pytesseract.image_to_string(denoised, config='--psm 6')
            if text1.strip():
                text_results.append(text1.strip())
            
            # Approach 2: High contrast
            _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            text2 = pytesseract.image_to_string(binary, config='--psm 6')
            if text2.strip():
                text_results.append(text2.strip())
            
            # Approach 3: Morphological operations for document text
            kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 1))
            morph = cv2.morphologyEx(gray, cv2.MORPH_CLOSE, kernel)
            text3 = pytesseract.image_to_string(morph, config='--psm 4')
            if text3.strip():
                text_results.append(text3.strip())
            
            # Approach 4: Edge enhancement
            edges = cv2.Canny(gray, 50, 150)
            dilated = cv2.dilate(edges, kernel, iterations=1)
            text4 = pytesseract.image_to_string(dilated, config='--psm 6')
            if text4.strip():
                text_results.append(text4.strip())
            
            # Select best result based on confidence and length
            best_text = ""
            best_score = 0
            
            for text in text_results:
                # Simple scoring based on length and character variety
                score = len(text) + len(set(text.lower())) * 2
                if score > best_score:
                    best_score = score
                    best_text = text
            
            return best_text
            
        except Exception as e:
            logger.warning(f"Enhanced OCR failed: {e}")
            return ""
    
    def _analyze_image_content(self, image: np.ndarray, file_path: Path) -> Dict[str, Any]:
        """Analyze image content using computer vision techniques"""
        try:
            analysis = {
                'description': '',
                'detected_objects': [],
                'scene_type': 'unknown',
                'text_regions': [],
                'visual_features': {}
            }
            
            # Basic scene analysis
            height, width = image.shape[:2]
            
            # Color analysis
            hsv = cv2.cvtColor(image, cv2.COLOR_BGR2HSV)
            dominant_colors = self._extract_dominant_colors(image)
            
            # Edge and contour analysis
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            edges = cv2.Canny(gray, 50, 150)
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            # Analyze image content type
            if self._is_document_image(image, contours):
                analysis['scene_type'] = 'document'
                analysis['description'] = self._generate_document_description(image, contours, file_path.name)
            elif self._is_diagram_or_chart(image, contours):
                analysis['scene_type'] = 'diagram_chart'
                analysis['description'] = self._generate_diagram_description(image, contours, file_path.name)
            elif self._is_screenshot_or_ui(image, contours):
                analysis['scene_type'] = 'screenshot_ui'
                analysis['description'] = self._generate_screenshot_description(image, contours, file_path.name)
            else:
                analysis['scene_type'] = 'general_image'
                analysis['description'] = self._generate_general_image_description(image, file_path.name)
            
            # Detect text regions
            analysis['text_regions'] = self._detect_text_regions(image)
            
            # Visual features
            analysis['visual_features'] = {
                'dominant_colors': dominant_colors,
                'edge_density': len(contours),
                'aspect_ratio': width / height,
                'brightness': cv2.mean(gray)[0] / 255.0,
                'contrast': gray.std() / 255.0
            }
            
            return analysis
            
        except Exception as e:
            logger.warning(f"Image content analysis failed: {e}")
            return {
                'description': f"Image file: {file_path.name}",
                'detected_objects': [],
                'scene_type': 'unknown',
                'text_regions': [],
                'visual_features': {}
            }
    
    def _analyze_document_structure_in_image(self, image: np.ndarray) -> Dict[str, Any]:
        """Analyze document structure in images (tables, forms, layouts)"""
        try:
            if not self.ocr_available:
                return {}
            
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            
            # Detect horizontal and vertical lines (table detection)
            horizontal_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (40, 1))
            vertical_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 40))
            
            # Detect horizontal lines
            horizontal_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, horizontal_kernel)
            # Detect vertical lines
            vertical_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, vertical_kernel)
            
            # Combine lines to detect tables
            table_mask = cv2.bitwise_or(horizontal_lines, vertical_lines)
            
            structural_elements = []
            
            if cv2.countNonZero(table_mask) > 1000:  # Threshold for table detection
                structural_elements.append("table")
                # Extract table content
                table_text = self._extract_table_content(image, table_mask)
                structural_text = f"Table Content:\n{table_text}"
            else:
                # Look for form-like structures
                form_elements = self._detect_form_elements(image)
                if form_elements:
                    structural_elements.extend(form_elements)
                    structural_text = f"Form Elements Detected: {', '.join(form_elements)}"
                else:
                    structural_text = ""
            
            return {
                'structural_elements': structural_elements,
                'structural_text': structural_text,
                'has_tables': 'table' in structural_elements,
                'has_forms': any('form' in elem for elem in structural_elements)
            }
            
        except Exception as e:
            logger.warning(f"Document structure analysis failed: {e}")
            return {}
    
    def _analyze_diagrams_and_charts(self, image: np.ndarray) -> Dict[str, Any]:
        """Analyze diagrams, flowcharts, and charts in images"""
        try:
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            edges = cv2.Canny(gray, 50, 150)
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            diagram_elements = []
            diagram_text_parts = []
            
            # Detect geometric shapes (flowchart elements)
            shapes = self._detect_geometric_shapes(contours)
            if shapes:
                diagram_elements.extend(shapes)
                diagram_text_parts.append(f"Geometric shapes detected: {', '.join(shapes)}")
            
            # Detect arrows and connectors
            arrows = self._detect_arrows_and_connectors(image, edges)
            if arrows:
                diagram_elements.append("connectors")
                diagram_text_parts.append(f"Flow connectors: {arrows} arrows/connectors detected")
            
            # Detect chart patterns (bars, lines, pie charts)
            chart_type = self._detect_chart_type(image, contours)
            if chart_type != 'unknown':
                diagram_elements.append(chart_type)
                diagram_text_parts.append(f"Chart type: {chart_type}")
            
            diagram_text = "\n".join(diagram_text_parts) if diagram_text_parts else ""
            
            return {
                'diagram_elements': diagram_elements,
                'diagram_text': diagram_text,
                'is_flowchart': 'rectangle' in shapes and 'connectors' in diagram_elements if 'shapes' in locals() else False,
                'is_chart': chart_type != 'unknown',
                'chart_type': chart_type if 'chart_type' in locals() else 'unknown'
            }
            
        except Exception as e:
            logger.warning(f"Diagram analysis failed: {e}")
            return {}
    
    def _combine_image_extractions(self, extracted_texts: List[str], methods: List[str]) -> str:
        """Intelligently combine text extracted from different methods"""
        if not extracted_texts:
            return ""
        
        # Remove duplicates while preserving order
        unique_texts = []
        seen = set()
        
        for text in extracted_texts:
            # Normalize text for comparison
            normalized = re.sub(r'\s+', ' ', text.strip().lower())
            if normalized and normalized not in seen:
                seen.add(normalized)
                unique_texts.append(text.strip())
        
        if len(unique_texts) == 1:
            return unique_texts[0]
        
        # Combine texts with method attribution
        combined_parts = []
        for i, (text, method) in enumerate(zip(unique_texts, methods)):
            if text:
                if method == "enhanced_ocr":
                    combined_parts.append(f"Text Content:\n{text}")
                elif method == "computer_vision":
                    combined_parts.append(f"Visual Analysis:\n{text}")
                elif method == "structure_analysis":
                    combined_parts.append(f"Document Structure:\n{text}")
                elif method == "diagram_analysis":
                    combined_parts.append(f"Diagram/Chart Analysis:\n{text}")
                else:
                    combined_parts.append(text)
        
        return "\n\n".join(combined_parts)
    
    # Helper methods for image analysis
    
    def _extract_dominant_colors(self, image: np.ndarray, k: int = 3) -> List[str]:
        """Extract dominant colors from image"""
        try:
            # Reshape image to be a list of pixels
            pixels = image.reshape(-1, 3)
            # Sample pixels for performance
            if len(pixels) > 10000:
                pixels = pixels[::len(pixels)//10000]
            
            # Simple dominant color detection (most frequent color ranges)
            colors = []
            mean_color = np.mean(pixels, axis=0)
            
            if mean_color[0] > 200 and mean_color[1] > 200 and mean_color[2] > 200:
                colors.append("light/white")
            elif mean_color[0] < 50 and mean_color[1] < 50 and mean_color[2] < 50:
                colors.append("dark/black")
            else:
                colors.append("mixed_colors")
            
            return colors
        except:
            return ["unknown"]
    
    def _is_document_image(self, image: np.ndarray, contours) -> bool:
        """Detect if image contains document-like content"""
        height, width = image.shape[:2]
        aspect_ratio = width / height
        
        # Document-like aspect ratios and text density
        if 0.7 <= aspect_ratio <= 1.5:  # Portrait or square documents
            # Check for horizontal text lines
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            horizontal_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (20, 1))
            horizontal_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, horizontal_kernel)
            return cv2.countNonZero(horizontal_lines) > width * 0.1
        
        return False
    
    def _is_diagram_or_chart(self, image: np.ndarray, contours) -> bool:
        """Detect if image contains diagrams or charts"""
        # Look for geometric shapes and patterns
        geometric_shapes = 0
        for contour in contours:
            if cv2.contourArea(contour) > 100:
                approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                if len(approx) in [3, 4, 5, 6, 8]:  # Triangle, rectangle, pentagon, hexagon, octagon
                    geometric_shapes += 1
        
        return geometric_shapes >= 2
    
    def _is_screenshot_or_ui(self, image: np.ndarray, contours) -> bool:
        """Detect if image is a screenshot or UI"""
        height, width = image.shape[:2]
        
        # Look for UI elements like buttons, windows, etc.
        rectangular_shapes = 0
        for contour in contours:
            if cv2.contourArea(contour) > 500:
                approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                if len(approx) == 4:  # Rectangle (common in UI)
                    rectangular_shapes += 1
        
        return rectangular_shapes >= 3
    
    def _generate_document_description(self, image: np.ndarray, contours, filename: str) -> str:
        """Generate description for document images"""
        height, width = image.shape[:2]
        return f"Document image '{filename}' ({width}x{height}px) containing text content and structured information. Suitable for OCR text extraction and document analysis."
    
    def _generate_diagram_description(self, image: np.ndarray, contours, filename: str) -> str:
        """Generate description for diagram/chart images"""
        return f"Diagram or chart image '{filename}' containing visual information, geometric shapes, and potentially data visualization elements. May include process flows, organizational charts, or data representations."
    
    def _generate_screenshot_description(self, image: np.ndarray, contours, filename: str) -> str:
        """Generate description for screenshot images"""
        return f"Screenshot or user interface image '{filename}' showing application interfaces, system windows, or digital content that may contain operational procedures or system information."
    
    def _generate_general_image_description(self, image: np.ndarray, filename: str) -> str:
        """Generate description for general images"""
        height, width = image.shape[:2]
        return f"Image file '{filename}' ({width}x{height}px) containing visual content that may include text, objects, or informational elements relevant to organizational knowledge."
    
    def _detect_text_regions(self, image: np.ndarray) -> List[Dict[str, Any]]:
        """Detect regions likely to contain text"""
        try:
            if not self.ocr_available:
                return []
            
            # Use pytesseract to get text regions
            boxes = pytesseract.image_to_boxes(image)
            regions = []
            
            for box in boxes.splitlines():
                parts = box.split(' ')
                if len(parts) >= 6:
                    char = parts[0]
                    x1, y1, x2, y2 = map(int, parts[1:5])
                    regions.append({
                        'character': char,
                        'bbox': [x1, y1, x2, y2],
                        'confidence': 0.8  # Default confidence
                    })
            
            return regions
        except:
            return []
    
    def _extract_table_content(self, image: np.ndarray, table_mask: np.ndarray) -> str:
        """Extract content from detected tables"""
        try:
            if not self.ocr_available:
                return "Table structure detected but OCR unavailable for content extraction"
            
            # Use OCR on the table region
            table_text = pytesseract.image_to_string(table_mask, config='--psm 6')
            return table_text.strip() if table_text.strip() else "Table detected but content unclear"
        except:
            return "Table structure detected"
    
    def _detect_form_elements(self, image: np.ndarray) -> List[str]:
        """Detect form elements like checkboxes, text fields"""
        elements = []
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        
        # Detect checkboxes (small squares)
        checkbox_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (10, 10))
        checkbox_detection = cv2.morphologyEx(gray, cv2.MORPH_OPEN, checkbox_kernel)
        if cv2.countNonZero(checkbox_detection) > 50:
            elements.append("checkboxes")
        
        # Detect horizontal lines (potential text fields)
        horizontal_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (30, 1))
        horizontal_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, horizontal_kernel)
        if cv2.countNonZero(horizontal_lines) > 100:
            elements.append("text_fields")
        
        return elements
    
    def _detect_geometric_shapes(self, contours) -> List[str]:
        """Detect geometric shapes in contours"""
        shapes = []
        shape_counts = {'triangle': 0, 'rectangle': 0, 'pentagon': 0, 'circle': 0}
        
        for contour in contours:
            if cv2.contourArea(contour) > 100:
                # Approximate contour
                approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                
                if len(approx) == 3:
                    shape_counts['triangle'] += 1
                elif len(approx) == 4:
                    shape_counts['rectangle'] += 1
                elif len(approx) == 5:
                    shape_counts['pentagon'] += 1
                elif len(approx) > 8:
                    shape_counts['circle'] += 1
        
        # Return shapes that were found
        for shape, count in shape_counts.items():
            if count > 0:
                shapes.append(shape)
        
        return shapes
    
    def _detect_arrows_and_connectors(self, image: np.ndarray, edges: np.ndarray) -> int:
        """Detect arrows and connectors in the image"""
        # Simple line detection
        lines = cv2.HoughLinesP(edges, 1, np.pi/180, threshold=50, minLineLength=30, maxLineGap=10)
        return len(lines) if lines is not None else 0
    
    def _detect_chart_type(self, image: np.ndarray, contours) -> str:
        """Detect type of chart in the image"""
        # Simple heuristics for chart detection
        height, width = image.shape[:2]
        
        # Look for rectangular regions (bar charts)
        rectangles = 0
        for contour in contours:
            if cv2.contourArea(contour) > 200:
                approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                if len(approx) == 4:
                    rectangles += 1
        
        if rectangles >= 3:
            return "bar_chart"
        
        # Look for circular regions (pie charts)
        circles = 0
        for contour in contours:
            if cv2.contourArea(contour) > 500:
                approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                if len(approx) > 8:
                    circles += 1
        
        if circles >= 1:
            return "pie_chart"
        
        return "unknown"
    
    def _detect_image_content_type(self, image: np.ndarray, text_content: str) -> str:
        """Detect the type of content in the image"""
        if not text_content:
            return "non_textual_image"
        
        text_lower = text_content.lower()
        
        if any(word in text_lower for word in ['procedure', 'step', 'instruction', 'process']):
            return "procedural_document"
        elif any(word in text_lower for word in ['policy', 'compliance', 'regulation', 'standard']):
            return "policy_document"
        elif any(word in text_lower for word in ['chart', 'graph', 'data', 'figure']):
            return "data_visualization"
        elif any(word in text_lower for word in ['form', 'application', 'request', 'checkbox']):
            return "form_document"
        elif any(word in text_lower for word in ['manual', 'guide', 'handbook', 'reference']):
            return "reference_document"
        else:
            return "general_document"
    
    def _identify_visual_elements(self, image: np.ndarray) -> Dict[str, Any]:
        """Identify visual elements in the image"""
        try:
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            edges = cv2.Canny(gray, 50, 150)
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            return {
                'total_contours': len(contours),
                'large_objects': len([c for c in contours if cv2.contourArea(c) > 1000]),
                'edge_density': cv2.countNonZero(edges) / (image.shape[0] * image.shape[1]),
                'has_geometric_shapes': len([c for c in contours if cv2.contourArea(c) > 100]) > 5
            }
        except:
            return {}
    
    def _basic_image_processing_fallback(self, file_path: Path) -> Dict[str, Any]:
        """Basic fallback processing for images when advanced methods fail"""
        try:
            image = cv2.imread(str(file_path))
            if image is None:
                raise ProcessingError("Unable to load image")
            
            height, width = image.shape[:2]
            
            # Try basic OCR if available
            text_content = ""
            if self.ocr_available:
                try:
                    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
                    text_content = pytesseract.image_to_string(gray)
                except:
                    pass
            
            return {
                'text': text_content or f"Image file: {file_path.name} - Advanced processing unavailable",
                'image_width': width,
                'image_height': height,
                'extraction_methods': ['fallback'],
                'sections': [],
                'processing_note': 'Basic fallback processing used'
            }
        except Exception as e:
            raise ProcessingError(f"Even fallback image processing failed: {e}")
    
    # Enhanced Video Processing Methods
    
    def _extract_video_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract comprehensive video metadata"""
        try:
            capture = cv2.VideoCapture(str(file_path))
            if not capture.isOpened():
                return {}
            
            fps = capture.get(cv2.CAP_PROP_FPS) or 0.0
            frame_count = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
            width = int(capture.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
            height = int(capture.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)
            duration = frame_count / fps if fps > 0 else 0
            
            capture.release()
            
            return {
                'fps': fps,
                'frame_count': frame_count,
                'width': width,
                'height': height,
                'duration': duration,
                'aspect_ratio': width / height if height > 0 else 0,
                'resolution': f"{width}x{height}",
                'estimated_size_mb': file_path.stat().st_size / (1024 * 1024)
            }
        except Exception as e:
            logger.warning(f"Failed to extract video metadata: {e}")
            return {}
    
    def _enhanced_audio_transcription(self, file_path: Path) -> Dict[str, Any]:
        """Enhanced audio transcription with better quality and metadata"""
        try:
            import ffmpeg
            
            # Create temporary audio file with optimized settings
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_audio:
                temp_audio_path = temp_audio.name
            
            try:
                # Enhanced audio extraction with noise reduction
                (
                    ffmpeg
                    .input(str(file_path))
                    .filter('highpass', f=200)  # Remove low-frequency noise
                    .filter('lowpass', f=8000)   # Remove high-frequency noise
                    .output(
                        temp_audio_path, 
                        acodec='pcm_s16le', 
                        ac=1, 
                        ar='16000',
                        af='volume=2.0'  # Boost volume
                    )
                    .overwrite_output()
                    .run(quiet=True, capture_stdout=True, capture_stderr=True)
                )
                
                # Transcribe with enhanced settings
                result = self.whisper_model.transcribe(
                    temp_audio_path,
                    fp16=False,
                    language=None,  # Auto-detect
                    task='transcribe',
                    verbose=False
                )
                
                transcript = (result.get('text') or '').strip()
                language = result.get('language', 'unknown')
                segments = result.get('segments', [])
                
                # Calculate confidence and quality metrics
                avg_confidence = 0.0
                if segments:
                    confidences = []
                    for segment in segments:
                        if 'avg_logprob' in segment:
                            # Convert log probability to confidence (rough approximation)
                            confidence = max(0.0, min(1.0, (segment['avg_logprob'] + 1.0)))
                            confidences.append(confidence)
                    avg_confidence = sum(confidences) / len(confidences) if confidences else 0.5
                
                return {
                    'transcript': transcript,
                    'language': language,
                    'segments': segments,
                    'metadata': {
                        'confidence': avg_confidence,
                        'segment_count': len(segments),
                        'audio_duration': result.get('duration', 0),
                        'processing_time': 0  # Could add timing if needed
                    }
                }
                
            finally:
                if os.path.exists(temp_audio_path):
                    os.unlink(temp_audio_path)
                    
        except Exception as e:
            logger.warning(f"Enhanced audio transcription failed: {e}")
            return {}
    
    def _intelligent_frame_analysis(self, file_path: Path) -> Dict[str, Any]:
        """Intelligent frame sampling and analysis"""
        try:
            capture = cv2.VideoCapture(str(file_path))
            if not capture.isOpened():
                return {}
            
            fps = capture.get(cv2.CAP_PROP_FPS) or 1.0
            total_frames = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
            duration = total_frames / fps
            
            # Intelligent sampling strategy
            if duration <= 30:  # Short video - sample every 2 seconds
                sample_interval = max(1, int(fps * 2))
                max_frames = 15
            elif duration <= 300:  # Medium video - sample every 10 seconds
                sample_interval = max(1, int(fps * 10))
                max_frames = 30
            else:  # Long video - sample every 30 seconds
                sample_interval = max(1, int(fps * 30))
                max_frames = 60
            
            frames_data = []
            text_content = []
            scene_descriptions = []
            frames_analyzed = 0
            
            frame_index = 0
            while frame_index < total_frames and frames_analyzed < max_frames:
                capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
                success, frame = capture.read()
                
                if not success or frame is None:
                    break
                
                # Analyze frame content
                frame_analysis = self._analyze_video_frame(frame, frame_index, fps)
                
                if frame_analysis.get('text'):
                    text_content.append(frame_analysis['text'])
                
                if frame_analysis.get('scene_description'):
                    scene_descriptions.append(f"Frame {frame_index} ({frame_index/fps:.1f}s): {frame_analysis['scene_description']}")
                
                frames_data.append(frame_analysis)
                frames_analyzed += 1
                frame_index += sample_interval
            
            capture.release()
            
            # Combine and deduplicate text content
            unique_text = self._deduplicate_video_text(text_content)
            combined_text = '\n'.join(unique_text) if unique_text else ""
            
            return {
                'visual_text': combined_text,
                'scene_analysis': '\n'.join(scene_descriptions),
                'metadata': {
                    'frames_analyzed': frames_analyzed,
                    'total_frames': total_frames,
                    'sampling_interval_seconds': sample_interval / fps,
                    'unique_text_blocks': len(unique_text),
                    'scene_changes_detected': len(scene_descriptions)
                }
            }
            
        except Exception as e:
            logger.warning(f"Intelligent frame analysis failed: {e}")
            return {}
    
    def _analyze_video_frame(self, frame: np.ndarray, frame_index: int, fps: float) -> Dict[str, Any]:
        """Analyze individual video frame for content"""
        try:
            analysis = {
                'frame_index': frame_index,
                'timestamp': frame_index / fps,
                'text': '',
                'scene_description': '',
                'visual_elements': {}
            }
            
            # Enhanced OCR for text extraction
            if self.ocr_available:
                text = self._enhanced_ocr_extraction(frame)
                if text and len(text.strip()) > 3:  # Filter very short text
                    analysis['text'] = text.strip()
            
            # Scene analysis
            scene_type = self._classify_video_scene(frame)
            analysis['scene_type'] = scene_type
            
            # Generate scene description based on type
            if scene_type == 'presentation':
                analysis['scene_description'] = "Presentation slide with text content"
            elif scene_type == 'document':
                analysis['scene_description'] = "Document or text-heavy content"
            elif scene_type == 'interface':
                analysis['scene_description'] = "Software interface or application screen"
            elif scene_type == 'diagram':
                analysis['scene_description'] = "Diagram, chart, or visual schematic"
            elif scene_type == 'procedural':
                analysis['scene_description'] = "Procedural content showing steps or actions"
            else:
                analysis['scene_description'] = "General video content"
            
            # Visual elements analysis
            analysis['visual_elements'] = self._identify_visual_elements(frame)
            
            return analysis
            
        except Exception as e:
            logger.warning(f"Frame analysis failed for frame {frame_index}: {e}")
            return {'frame_index': frame_index, 'timestamp': frame_index / fps}
    
    def _classify_video_scene(self, frame: np.ndarray) -> str:
        """Classify the type of scene in a video frame"""
        try:
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            height, width = gray.shape
            
            # Text density analysis
            if self.ocr_available:
                try:
                    data = pytesseract.image_to_data(frame, output_type=pytesseract.Output.DICT)
                    text_regions = len([conf for conf in data['conf'] if int(conf) > 30])
                    text_density = text_regions / (width * height / 10000)  # Normalize by area
                    
                    if text_density > 5:
                        return 'document'
                    elif text_density > 2:
                        return 'presentation'
                except:
                    text_density = 0
            else:
                text_density = 0
            
            # Edge and contour analysis
            edges = cv2.Canny(gray, 50, 150)
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            # Rectangle detection (UI elements)
            rectangles = 0
            for contour in contours:
                if cv2.contourArea(contour) > 500:
                    approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                    if len(approx) == 4:
                        rectangles += 1
            
            if rectangles > 5:
                return 'interface'
            
            # Geometric shapes (diagrams)
            geometric_shapes = 0
            for contour in contours:
                if cv2.contourArea(contour) > 200:
                    approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                    if len(approx) in [3, 4, 5, 6, 8]:
                        geometric_shapes += 1
            
            if geometric_shapes > 3:
                return 'diagram'
            
            # Check for procedural content indicators
            if text_density > 0.5 and rectangles > 1:
                return 'procedural'
            
            return 'general'
            
        except Exception as e:
            logger.warning(f"Scene classification failed: {e}")
            return 'general'
    
    def _detect_procedural_content(self, file_path: Path, extracted_content: Dict[str, Any]) -> Dict[str, Any]:
        """Detect procedural content in videos (tutorials, instructions, etc.)"""
        try:
            # Analyze extracted content for procedural indicators
            audio_text = extracted_content.get('audio_transcript', '')
            visual_text = extracted_content.get('visual_text', '')
            
            procedural_indicators = [
                'step', 'first', 'next', 'then', 'now', 'click', 'select', 'choose',
                'open', 'close', 'press', 'type', 'enter', 'drag', 'drop',
                'procedure', 'tutorial', 'how to', 'instruction', 'guide'
            ]
            
            combined_text = f"{audio_text} {visual_text}".lower()
            indicator_count = sum(1 for indicator in procedural_indicators if indicator in combined_text)
            
            if indicator_count >= 3:  # Threshold for procedural content
                # Extract step-by-step information
                procedural_text = self._extract_procedural_steps(audio_text, visual_text)
                
                return {
                    'procedural_text': procedural_text,
                    'metadata': {
                        'procedural_indicator_count': indicator_count,
                        'is_procedural': True,
                        'confidence': min(1.0, indicator_count / 10.0)
                    }
                }
            
            return {'metadata': {'is_procedural': False}}
            
        except Exception as e:
            logger.warning(f"Procedural content detection failed: {e}")
            return {}
    
    def _analyze_video_scenes(self, file_path: Path) -> Dict[str, Any]:
        """Analyze video for scene changes and content categorization"""
        try:
            capture = cv2.VideoCapture(str(file_path))
            if not capture.isOpened():
                return {}
            
            fps = capture.get(cv2.CAP_PROP_FPS) or 1.0
            total_frames = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
            
            # Sample frames for scene analysis
            scene_changes = []
            prev_frame_features = None
            scenes = []
            
            # Sample every 5 seconds for scene analysis
            sample_interval = max(1, int(fps * 5))
            
            frame_index = 0
            scene_count = 0
            
            while frame_index < total_frames and scene_count < 20:  # Limit scenes
                capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
                success, frame = capture.read()
                
                if not success or frame is None:
                    break
                
                # Calculate frame features for scene change detection
                frame_features = self._calculate_frame_features(frame)
                
                if prev_frame_features is not None:
                    # Simple scene change detection
                    difference = self._calculate_frame_difference(frame_features, prev_frame_features)
                    if difference > 0.3:  # Threshold for scene change
                        scene_changes.append(frame_index / fps)
                        scene_type = self._classify_video_scene(frame)
                        scenes.append(f"Scene {scene_count + 1} at {frame_index/fps:.1f}s: {scene_type}")
                        scene_count += 1
                
                prev_frame_features = frame_features
                frame_index += sample_interval
            
            capture.release()
            
            return {
                'scene_descriptions': '\n'.join(scenes),
                'metadata': {
                    'scene_count': len(scenes),
                    'scene_changes': scene_changes,
                    'average_scene_length': (total_frames / fps) / max(len(scenes), 1)
                }
            }
            
        except Exception as e:
            logger.warning(f"Scene analysis failed: {e}")
            return {}
    
    def _combine_video_extractions(self, extracted_content: Dict[str, Any], sources: List[str]) -> str:
        """Intelligently combine content extracted from video"""
        parts = []
        
        # Audio transcript (usually most important)
        if extracted_content.get('audio_transcript'):
            parts.append(f"Audio Transcript:\n{extracted_content['audio_transcript']}")
        
        # Visual text content
        if extracted_content.get('visual_text'):
            parts.append(f"On-Screen Text:\n{extracted_content['visual_text']}")
        
        # Scene analysis
        if extracted_content.get('scene_analysis'):
            parts.append(f"Visual Scene Analysis:\n{extracted_content['scene_analysis']}")
        
        # Procedural content
        if extracted_content.get('procedural_content'):
            parts.append(f"Procedural Information:\n{extracted_content['procedural_content']}")
        
        return '\n\n'.join(parts) if parts else ""
    
    def _basic_video_processing_fallback(self, file_path: Path) -> Dict[str, Any]:
        """Basic fallback video processing when advanced methods fail"""
        try:
            # Try the original frame OCR method
            ocr_result = self._extract_text_from_video_frames(
                file_path=file_path,
                interval_seconds=10,  # Less frequent sampling for fallback
                max_frames=10
            )
            
            text_content = ocr_result.get('text', '').strip()
            if not text_content:
                text_content = f"Video file: {file_path.name} - Content extraction unavailable"
            
            return {
                'text': text_content,
                'frames_analyzed': ocr_result.get('frames_analyzed', 0),
                'processing_note': 'Basic fallback processing used'
            }
            
        except Exception as e:
            logger.warning(f"Even basic video processing failed: {e}")
            return {
                'text': f"Video file: {file_path.name} - Processing failed",
                'frames_analyzed': 0,
                'processing_note': 'All video processing methods failed'
            }
    
    # Helper methods for video processing
    
    def _deduplicate_video_text(self, text_list: List[str]) -> List[str]:
        """Remove duplicate text content from video frames"""
        unique_texts = []
        seen_normalized = set()
        
        for text in text_list:
            if not text or len(text.strip()) < 5:
                continue
            
            # Normalize for comparison
            normalized = re.sub(r'\s+', ' ', text.strip().lower())
            normalized = re.sub(r'[^\w\s]', '', normalized)  # Remove punctuation
            
            if normalized not in seen_normalized:
                seen_normalized.add(normalized)
                unique_texts.append(text.strip())
        
        return unique_texts
    
    def _extract_procedural_steps(self, audio_text: str, visual_text: str) -> str:
        """Extract procedural steps from audio and visual content"""
        combined_text = f"{audio_text}\n{visual_text}"
        
        # Look for step patterns
        step_patterns = [
            r'step (\d+)[:\.](.+?)(?=step \d+|$)',
            r'(\d+)\.(.+?)(?=\d+\.|$)',
            r'(first|second|third|fourth|fifth|next|then|finally)[,:](.+?)(?=first|second|third|fourth|fifth|next|then|finally|$)'
        ]
        
        extracted_steps = []
        for pattern in step_patterns:
            matches = re.finditer(pattern, combined_text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                if len(match.groups()) >= 2:
                    step_content = match.group(2).strip()
                    if len(step_content) > 10:
                        extracted_steps.append(f"Step: {step_content}")
        
        return '\n'.join(extracted_steps) if extracted_steps else ""
    
    def _calculate_frame_features(self, frame: np.ndarray) -> np.ndarray:
        """Calculate features for frame comparison"""
        try:
            # Convert to grayscale and resize for speed
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            small = cv2.resize(gray, (64, 64))
            
            # Calculate histogram as feature
            hist = cv2.calcHist([small], [0], None, [32], [0, 256])
            return hist.flatten()
        except:
            return np.zeros(32)
    
    def _calculate_frame_difference(self, features1: np.ndarray, features2: np.ndarray) -> float:
        """Calculate difference between frame features"""
        try:
            # Normalize features
            norm1 = features1 / (np.sum(features1) + 1e-6)
            norm2 = features2 / (np.sum(features2) + 1e-6)
            
            # Calculate histogram difference
            diff = np.sum(np.abs(norm1 - norm2))
            return diff
        except:
            return 0.0
    
    def _detect_video_content_type(self, text_content: str, metadata: Dict[str, Any]) -> str:
        """Detect the type of content in the video"""
        if not text_content:
            return "non_textual_video"
        
        text_lower = text_content.lower()
        
        if any(word in text_lower for word in ['tutorial', 'how to', 'step', 'instruction', 'guide']):
            return "instructional_video"
        elif any(word in text_lower for word in ['presentation', 'slide', 'agenda', 'overview']):
            return "presentation_video"
        elif any(word in text_lower for word in ['procedure', 'process', 'workflow', 'protocol']):
            return "procedural_video"
        elif any(word in text_lower for word in ['meeting', 'discussion', 'review', 'interview']):
            return "meeting_video"
        elif any(word in text_lower for word in ['training', 'course', 'lesson', 'education']):
            return "training_video"
        elif any(word in text_lower for word in ['demo', 'demonstration', 'example', 'showcase']):
            return "demonstration_video"
        else:
            return "general_video"
    
    def _calculate_video_extraction_confidence(self, extracted_content: Dict[str, Any], sources: List[str]) -> float:
        """Calculate confidence score for video knowledge extraction"""
        confidence_factors = []
        
        # Audio transcript confidence
        if extracted_content.get('audio_transcript'):
            audio_length = len(extracted_content['audio_transcript'])
            audio_confidence = min(1.0, audio_length / 500)  # Normalize by expected length
            confidence_factors.append(audio_confidence * 0.4)  # 40% weight
        
        # Visual text confidence
        if extracted_content.get('visual_text'):
            visual_length = len(extracted_content['visual_text'])
            visual_confidence = min(1.0, visual_length / 200)
            confidence_factors.append(visual_confidence * 0.3)  # 30% weight
        
        # Scene analysis confidence
        if extracted_content.get('scene_analysis'):
            scene_confidence = 0.8  # Fixed confidence for scene analysis
            confidence_factors.append(scene_confidence * 0.2)  # 20% weight
        
        # Procedural content confidence
        if extracted_content.get('procedural_content'):
            procedural_confidence = 0.9  # High confidence for procedural content
            confidence_factors.append(procedural_confidence * 0.1)  # 10% weight
        
        # Calculate overall confidence
        if confidence_factors:
            return sum(confidence_factors) / len(confidence_factors)
        else:
            return 0.1  # Very low confidence if nothing extracted
