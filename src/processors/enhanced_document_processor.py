"""
Enhanced Document Processor for EXPLAINIUM

This module provides advanced document processing capabilities supporting multiple
input formats including text documents, images, videos, audio, spreadsheets, and
presentations with comprehensive knowledge extraction.
"""

import os
import io
import json
import logging
from typing import Dict, Any, List, Optional, Tuple, Union
from pathlib import Path
import tempfile
import shutil
from datetime import datetime

# Core processing libraries
import requests
import ffmpeg
import whisper
import cv2
import numpy as np
from PIL import Image, ImageEnhance
import pytesseract
import pandas as pd
import openpyxl
from pptx import Presentation
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import easyocr

# AI and NLP libraries
import spacy
from transformers import pipeline
import torch

# Internal imports
from src.logging_config import get_logger, log_processing_step, log_error
from src.config import config_manager
from src.exceptions import ProcessingError, AIError, ServiceUnavailableError
from src.ai.enhanced_knowledge_extractor import EnhancedKnowledgeExtractor
from src.ai.ner_extractor import ner_extractor
from src.ai.classifier import classifier
from src.ai.keyphrase_extractor import keyphrase_extractor

logger = get_logger(__name__)

# Configuration
TIKA_SERVER_URL = config_manager.get_tika_url()
SUPPORTED_IMAGE_FORMATS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
SUPPORTED_VIDEO_FORMATS = {'.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm'}
SUPPORTED_AUDIO_FORMATS = {'.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a'}
SUPPORTED_DOCUMENT_FORMATS = {'.pdf', '.doc', '.docx', '.txt', '.rtf'}
SUPPORTED_PRESENTATION_FORMATS = {'.ppt', '.pptx'}
SUPPORTED_SPREADSHEET_FORMATS = {'.xls', '.xlsx', '.csv'}

class EnhancedDocumentProcessor:
    """Enhanced document processor with multi-format support"""
    
    def __init__(self):
        self.knowledge_extractor = EnhancedKnowledgeExtractor()
        self.ocr_reader = None
        self.whisper_model = None
        self._initialize_models()
    
    def _initialize_models(self):
        """Initialize AI models and OCR engines"""
        try:
            # Initialize OCR reader
            self.ocr_reader = easyocr.Reader(['en'], gpu=torch.cuda.is_available())
            
            # Initialize Whisper model for audio transcription
            self.whisper_model = whisper.load_model("base")
            
            logger.info("Enhanced document processor models initialized successfully")
        except Exception as e:
            logger.error(f"Failed to initialize models: {e}")
            self.ocr_reader = None
            self.whisper_model = None
    
    def process_document(self, file_path: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        Main document processing method supporting multiple formats
        
        Args:
            file_path: Path to the document file
            metadata: Document metadata including document_id, filename, etc.
            
        Returns:
            Dictionary containing extracted text, knowledge, and processing results
        """
        logger.info(f"Starting enhanced processing for: {file_path}")
        
        # Initialize results structure
        results = {
            'extracted_text': '',
            'document_sections': {},
            'extracted_knowledge': {},
            'processing_metadata': {
                'file_type': self._get_file_type(file_path),
                'file_size': os.path.getsize(file_path),
                'processing_start': datetime.utcnow().isoformat(),
                'processing_methods': []
            },
            'confidence_scores': {},
            'processing_errors': []
        }
        
        try:
            # Determine file type and process accordingly
            file_type = self._get_file_type(file_path)
            
            if file_type == 'image':
                results = self._process_image(file_path, results, metadata)
            elif file_type == 'video':
                results = self._process_video(file_path, results, metadata)
            elif file_type == 'audio':
                results = self._process_audio(file_path, results, metadata)
            elif file_type == 'document':
                results = self._process_document_file(file_path, results, metadata)
            elif file_type == 'presentation':
                results = self._process_presentation(file_path, results, metadata)
            elif file_type == 'spreadsheet':
                results = self._process_spreadsheet(file_path, results, metadata)
            else:
                results = self._process_generic_file(file_path, results, metadata)
            
            # Extract comprehensive knowledge if we have text
            if results['extracted_text']:
                knowledge_results = self.knowledge_extractor.extract_comprehensive_knowledge(
                    results['extracted_text'], 
                    metadata
                )
                results['extracted_knowledge'] = knowledge_results
                results['confidence_scores'].update(knowledge_results.get('confidence_scores', {}))
            
            # Add processing completion metadata
            results['processing_metadata']['processing_end'] = datetime.utcnow().isoformat()
            results['processing_metadata']['success'] = True
            
            logger.info(f"Enhanced processing completed successfully for: {file_path}")
            
        except Exception as e:
            logger.error(f"Error in enhanced document processing: {e}")
            results['processing_errors'].append(str(e))
            results['processing_metadata']['success'] = False
            results['processing_metadata']['processing_end'] = datetime.utcnow().isoformat()
        
        return results
    
    def _get_file_type(self, file_path: str) -> str:
        """Determine file type based on extension"""
        ext = Path(file_path).suffix.lower()
        
        if ext in SUPPORTED_IMAGE_FORMATS:
            return 'image'
        elif ext in SUPPORTED_VIDEO_FORMATS:
            return 'video'
        elif ext in SUPPORTED_AUDIO_FORMATS:
            return 'audio'
        elif ext in SUPPORTED_DOCUMENT_FORMATS:
            return 'document'
        elif ext in SUPPORTED_PRESENTATION_FORMATS:
            return 'presentation'
        elif ext in SUPPORTED_SPREADSHEET_FORMATS:
            return 'spreadsheet'
        else:
            return 'unknown'
    
    def _process_image(self, file_path: str, results: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process image files with OCR and visual analysis"""
        logger.info(f"Processing image: {file_path}")
        results['processing_metadata']['processing_methods'].append('image_ocr')
        
        try:
            # Load and preprocess image
            image = Image.open(file_path)
            
            # Enhance image for better OCR
            enhanced_image = self._enhance_image_for_ocr(image)
            
            # Extract text using multiple OCR methods
            ocr_results = self._extract_text_from_image(enhanced_image, file_path)
            results['extracted_text'] = ocr_results['text']
            results['processing_metadata']['ocr_confidence'] = ocr_results['confidence']
            
            # Analyze image content
            image_analysis = self._analyze_image_content(image)
            results['document_sections']['image_analysis'] = image_analysis
            
            # Extract visual elements (charts, diagrams, etc.)
            visual_elements = self._extract_visual_elements(image)
            if visual_elements:
                results['document_sections']['visual_elements'] = visual_elements
            
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            results['processing_errors'].append(f"Image processing error: {str(e)}")
        
        return results
    
    def _enhance_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Enhance image quality for better OCR results"""
        try:
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Enhance contrast and sharpness
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(1.2)
            
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(1.1)
            
            # Resize if image is too small
            width, height = image.size
            if width < 1000 or height < 1000:
                scale_factor = max(1000 / width, 1000 / height)
                new_width = int(width * scale_factor)
                new_height = int(height * scale_factor)
                image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            return image
        except Exception as e:
            logger.error(f"Error enhancing image: {e}")
            return image
    
    def _extract_text_from_image(self, image: Image.Image, file_path: str) -> Dict[str, Any]:
        """Extract text from image using multiple OCR methods"""
        extracted_texts = []
        confidences = []
        
        try:
            # Method 1: Tesseract OCR
            tesseract_text = pytesseract.image_to_string(image, config='--psm 6')
            if tesseract_text.strip():
                extracted_texts.append(tesseract_text)
                # Get confidence from Tesseract
                tesseract_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                tesseract_conf = np.mean([int(conf) for conf in tesseract_data['conf'] if int(conf) > 0])
                confidences.append(tesseract_conf / 100.0)
        except Exception as e:
            logger.warning(f"Tesseract OCR failed: {e}")
        
        try:
            # Method 2: EasyOCR
            if self.ocr_reader:
                easyocr_results = self.ocr_reader.readtext(np.array(image))
                easyocr_text = ' '.join([result[1] for result in easyocr_results])
                if easyocr_text.strip():
                    extracted_texts.append(easyocr_text)
                    easyocr_conf = np.mean([result[2] for result in easyocr_results])
                    confidences.append(easyocr_conf)
        except Exception as e:
            logger.warning(f"EasyOCR failed: {e}")
        
        # Combine results and select best
        if extracted_texts:
            # Use the text with highest confidence
            best_idx = np.argmax(confidences) if confidences else 0
            best_text = extracted_texts[best_idx]
            best_confidence = confidences[best_idx] if confidences else 0.5
            
            return {
                'text': best_text,
                'confidence': best_confidence,
                'methods_used': len(extracted_texts)
            }
        else:
            return {
                'text': '',
                'confidence': 0.0,
                'methods_used': 0
            }
    
    def _analyze_image_content(self, image: Image.Image) -> Dict[str, Any]:
        """Analyze image content for visual elements and structure"""
        analysis = {
            'dimensions': image.size,
            'mode': image.mode,
            'format': image.format,
            'has_text_regions': False,
            'has_diagrams': False,
            'has_charts': False,
            'color_analysis': {}
        }
        
        try:
            # Convert to OpenCV format for analysis
            cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
            
            # Detect text regions
            text_regions = self._detect_text_regions(cv_image)
            analysis['has_text_regions'] = len(text_regions) > 0
            analysis['text_region_count'] = len(text_regions)
            
            # Detect shapes and diagrams
            shapes = self._detect_shapes(cv_image)
            analysis['has_diagrams'] = len(shapes) > 5  # Threshold for diagram detection
            analysis['shape_count'] = len(shapes)
            
            # Basic color analysis
            analysis['color_analysis'] = self._analyze_colors(image)
            
        except Exception as e:
            logger.warning(f"Image content analysis failed: {e}")
        
        return analysis
    
    def _detect_text_regions(self, cv_image: np.ndarray) -> List[Tuple[int, int, int, int]]:
        """Detect text regions in image using OpenCV"""
        try:
            # Convert to grayscale
            gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)
            
            # Apply morphological operations to detect text regions
            kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (18, 18))
            connected = cv2.morphologyEx(gray, cv2.MORPH_CLOSE, kernel)
            
            # Find contours
            contours, _ = cv2.findContours(connected, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            # Filter contours by area and aspect ratio
            text_regions = []
            for contour in contours:
                x, y, w, h = cv2.boundingRect(contour)
                area = w * h
                aspect_ratio = w / h if h > 0 else 0
                
                # Filter based on size and aspect ratio
                if area > 100 and 0.1 < aspect_ratio < 10:
                    text_regions.append((x, y, w, h))
            
            return text_regions
        except Exception as e:
            logger.warning(f"Text region detection failed: {e}")
            return []
    
    def _detect_shapes(self, cv_image: np.ndarray) -> List[Dict[str, Any]]:
        """Detect geometric shapes in image"""
        try:
            gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)
            edges = cv2.Canny(gray, 50, 150)
            
            # Find contours
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            shapes = []
            for contour in contours:
                # Approximate contour
                epsilon = 0.02 * cv2.arcLength(contour, True)
                approx = cv2.approxPolyDP(contour, epsilon, True)
                
                # Classify shape based on number of vertices
                vertices = len(approx)
                area = cv2.contourArea(contour)
                
                if area > 100:  # Filter small shapes
                    shape_info = {
                        'vertices': vertices,
                        'area': area,
                        'perimeter': cv2.arcLength(contour, True)
                    }
                    
                    if vertices == 3:
                        shape_info['type'] = 'triangle'
                    elif vertices == 4:
                        shape_info['type'] = 'rectangle'
                    elif vertices > 8:
                        shape_info['type'] = 'circle'
                    else:
                        shape_info['type'] = 'polygon'
                    
                    shapes.append(shape_info)
            
            return shapes
        except Exception as e:
            logger.warning(f"Shape detection failed: {e}")
            return []
    
    def _analyze_colors(self, image: Image.Image) -> Dict[str, Any]:
        """Analyze color distribution in image"""
        try:
            # Get dominant colors
            image_array = np.array(image)
            pixels = image_array.reshape(-1, 3)
            
            # Calculate color statistics
            mean_color = np.mean(pixels, axis=0)
            std_color = np.std(pixels, axis=0)
            
            return {
                'mean_rgb': mean_color.tolist(),
                'std_rgb': std_color.tolist(),
                'is_grayscale': len(np.unique(pixels)) < 256,
                'brightness': np.mean(mean_color)
            }
        except Exception as e:
            logger.warning(f"Color analysis failed: {e}")
            return {}
    
    def _extract_visual_elements(self, image: Image.Image) -> Dict[str, Any]:
        """Extract visual elements like charts, tables, diagrams"""
        elements = {
            'tables': [],
            'charts': [],
            'diagrams': [],
            'flowcharts': []
        }
        
        try:
            # Convert to OpenCV format
            cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
            
            # Detect table-like structures
            tables = self._detect_tables(cv_image)
            elements['tables'] = tables
            
            # Detect chart-like structures
            charts = self._detect_charts(cv_image)
            elements['charts'] = charts
            
        except Exception as e:
            logger.warning(f"Visual element extraction failed: {e}")
        
        return elements if any(elements.values()) else {}
    
    def _detect_tables(self, cv_image: np.ndarray) -> List[Dict[str, Any]]:
        """Detect table structures in image"""
        try:
            gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)
            
            # Detect horizontal and vertical lines
            horizontal_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (40, 1))
            vertical_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 40))
            
            horizontal_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, horizontal_kernel)
            vertical_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, vertical_kernel)
            
            # Combine lines
            table_mask = cv2.add(horizontal_lines, vertical_lines)
            
            # Find contours of potential tables
            contours, _ = cv2.findContours(table_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            tables = []
            for contour in contours:
                x, y, w, h = cv2.boundingRect(contour)
                area = w * h
                
                # Filter based on size and aspect ratio
                if area > 1000 and w > 100 and h > 100:
                    tables.append({
                        'bounds': (x, y, w, h),
                        'area': area,
                        'confidence': 0.7  # Basic confidence score
                    })
            
            return tables
        except Exception as e:
            logger.warning(f"Table detection failed: {e}")
            return []
    
    def _detect_charts(self, cv_image: np.ndarray) -> List[Dict[str, Any]]:
        """Detect chart structures in image"""
        try:
            gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)
            
            # Look for circular shapes (pie charts)
            circles = cv2.HoughCircles(gray, cv2.HOUGH_GRADIENT, 1, 20,
                                     param1=50, param2=30, minRadius=20, maxRadius=200)
            
            charts = []
            if circles is not None:
                circles = np.round(circles[0, :]).astype("int")
                for (x, y, r) in circles:
                    charts.append({
                        'type': 'pie_chart',
                        'center': (x, y),
                        'radius': r,
                        'confidence': 0.6
                    })
            
            # Look for bar chart patterns (rectangular shapes in sequence)
            contours, _ = cv2.findContours(gray, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            rectangles = []
            
            for contour in contours:
                approx = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
                if len(approx) == 4:  # Rectangle
                    x, y, w, h = cv2.boundingRect(contour)
                    rectangles.append((x, y, w, h))
            
            # Group rectangles that might form a bar chart
            if len(rectangles) > 3:  # Need multiple bars
                charts.append({
                    'type': 'bar_chart',
                    'elements': len(rectangles),
                    'confidence': 0.5
                })
            
            return charts
        except Exception as e:
            logger.warning(f"Chart detection failed: {e}")
            return []
    
    def _process_video(self, file_path: str, results: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process video files with audio extraction and frame analysis"""
        logger.info(f"Processing video: {file_path}")
        results['processing_metadata']['processing_methods'].append('video_processing')
        
        try:
            # Extract audio and transcribe
            audio_results = self._extract_and_transcribe_audio(file_path)
            if audio_results['text']:
                results['extracted_text'] = audio_results['text']
                results['document_sections']['audio_transcription'] = audio_results
            
            # Extract key frames and analyze
            frame_analysis = self._analyze_video_frames(file_path)
            if frame_analysis:
                results['document_sections']['frame_analysis'] = frame_analysis
                
                # Extract text from key frames
                frame_text = self._extract_text_from_frames(file_path, frame_analysis['key_frames'])
                if frame_text:
                    results['extracted_text'] += '\n\n' + frame_text
                    results['document_sections']['frame_text'] = frame_text
            
            # Extract video metadata
            video_metadata = self._extract_video_metadata(file_path)
            results['document_sections']['video_metadata'] = video_metadata
            
        except Exception as e:
            logger.error(f"Error processing video {file_path}: {e}")
            results['processing_errors'].append(f"Video processing error: {str(e)}")
        
        return results
    
    def _extract_and_transcribe_audio(self, video_path: str) -> Dict[str, Any]:
        """Extract audio from video and transcribe to text"""
        audio_results = {'text': '', 'confidence': 0.0, 'duration': 0.0}
        
        try:
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_audio:
                # Extract audio using ffmpeg
                stream = ffmpeg.input(video_path)
                stream = ffmpeg.output(stream, temp_audio.name, acodec='pcm_s16le', ac=1, ar='16000')
                ffmpeg.run(stream, overwrite_output=True, quiet=True)
                
                # Transcribe using Whisper
                if self.whisper_model:
                    result = self.whisper_model.transcribe(temp_audio.name)
                    audio_results['text'] = result['text']
                    
                    # Calculate average confidence
                    if 'segments' in result:
                        confidences = [seg.get('no_speech_prob', 0.5) for seg in result['segments']]
                        audio_results['confidence'] = 1.0 - np.mean(confidences) if confidences else 0.5
                    else:
                        audio_results['confidence'] = 0.7  # Default confidence
                    
                    # Get duration
                    audio_results['duration'] = result.get('duration', 0.0)
                
                # Clean up temporary file
                os.unlink(temp_audio.name)
                
        except Exception as e:
            logger.error(f"Audio extraction and transcription failed: {e}")
        
        return audio_results
    
    def _analyze_video_frames(self, video_path: str) -> Dict[str, Any]:
        """Analyze video frames to extract key frames and visual information"""
        try:
            cap = cv2.VideoCapture(video_path)
            
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            duration = frame_count / fps if fps > 0 else 0
            
            # Extract key frames (every 30 seconds or significant scene changes)
            key_frames = []
            frame_interval = max(1, int(fps * 30))  # Every 30 seconds
            
            for i in range(0, frame_count, frame_interval):
                cap.set(cv2.CAP_PROP_POS_FRAMES, i)
                ret, frame = cap.read()
                if ret:
                    timestamp = i / fps if fps > 0 else 0
                    key_frames.append({
                        'frame_number': i,
                        'timestamp': timestamp,
                        'has_text': self._frame_has_text(frame)
                    })
            
            cap.release()
            
            return {
                'total_frames': frame_count,
                'fps': fps,
                'duration': duration,
                'key_frames': key_frames
            }
            
        except Exception as e:
            logger.error(f"Video frame analysis failed: {e}")
            return {}
    
    def _frame_has_text(self, frame: np.ndarray) -> bool:
        """Check if a video frame contains readable text"""
        try:
            # Convert to PIL Image
            rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            pil_image = Image.fromarray(rgb_frame)
            
            # Quick OCR check
            text = pytesseract.image_to_string(pil_image, config='--psm 8')
            return len(text.strip()) > 10  # Threshold for meaningful text
            
        except Exception:
            return False
    
    def _extract_text_from_frames(self, video_path: str, key_frames: List[Dict[str, Any]]) -> str:
        """Extract text from key video frames"""
        extracted_texts = []
        
        try:
            cap = cv2.VideoCapture(video_path)
            
            for frame_info in key_frames:
                if frame_info.get('has_text', False):
                    cap.set(cv2.CAP_PROP_POS_FRAMES, frame_info['frame_number'])
                    ret, frame = cap.read()
                    
                    if ret:
                        # Convert to PIL Image
                        rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                        pil_image = Image.fromarray(rgb_frame)
                        
                        # Extract text
                        ocr_result = self._extract_text_from_image(pil_image, 'video_frame')
                        if ocr_result['text'].strip():
                            timestamp = frame_info.get('timestamp', 0)
                            extracted_texts.append(f"[Frame at {timestamp:.1f}s]: {ocr_result['text']}")
            
            cap.release()
            
        except Exception as e:
            logger.error(f"Frame text extraction failed: {e}")
        
        return '\n'.join(extracted_texts)
    
    def _extract_video_metadata(self, video_path: str) -> Dict[str, Any]:
        """Extract metadata from video file"""
        try:
            cap = cv2.VideoCapture(video_path)
            
            metadata = {
                'width': int(cap.get(cv2.CAP_PROP_FRAME_WIDTH)),
                'height': int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT)),
                'fps': cap.get(cv2.CAP_PROP_FPS),
                'frame_count': int(cap.get(cv2.CAP_PROP_FRAME_COUNT)),
                'codec': int(cap.get(cv2.CAP_PROP_FOURCC)),
                'file_size': os.path.getsize(video_path)
            }
            
            # Calculate duration
            if metadata['fps'] > 0:
                metadata['duration'] = metadata['frame_count'] / metadata['fps']
            
            cap.release()
            return metadata
            
        except Exception as e:
            logger.error(f"Video metadata extraction failed: {e}")
            return {}
    
    def _process_audio(self, file_path: str, results: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process audio files with transcription"""
        logger.info(f"Processing audio: {file_path}")
        results['processing_metadata']['processing_methods'].append('audio_transcription')
        
        try:
            if self.whisper_model:
                # Transcribe audio directly
                result = self.whisper_model.transcribe(file_path)
                results['extracted_text'] = result['text']
                
                # Add transcription details
                transcription_info = {
                    'text': result['text'],
                    'language': result.get('language', 'unknown'),
                    'duration': result.get('duration', 0.0),
                    'segments': len(result.get('segments', []))
                }
                
                # Calculate confidence
                if 'segments' in result:
                    confidences = [1.0 - seg.get('no_speech_prob', 0.5) for seg in result['segments']]
                    transcription_info['confidence'] = np.mean(confidences) if confidences else 0.5
                
                results['document_sections']['audio_transcription'] = transcription_info
                results['confidence_scores']['audio_transcription'] = transcription_info.get('confidence', 0.5)
            
        except Exception as e:
            logger.error(f"Error processing audio {file_path}: {e}")
            results['processing_errors'].append(f"Audio processing error: {str(e)}")
        
        return results
    
    def _process_document_file(self, file_path: str, results: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process document files (PDF, DOC, DOCX, TXT)"""
        logger.info(f"Processing document: {file_path}")
        results['processing_metadata']['processing_methods'].append('document_extraction')
        
        ext = Path(file_path).suffix.lower()
        
        try:
            if ext == '.pdf':
                results = self._process_pdf(file_path, results)
            elif ext in ['.doc', '.docx']:
                results = self._process_word_document(file_path, results)
            elif ext == '.txt':
                results = self._process_text_file(file_path, results)
            elif ext == '.rtf':
                results = self._process_rtf_file(file_path, results)
            else:
                # Fallback to Tika if available
                results = self._process_with_tika(file_path, results)
                
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            results['processing_errors'].append(f"Document processing error: {str(e)}")
        
        return results
    
    def _process_pdf(self, file_path: str, results: Dict[str, Any]) -> Dict[str, Any]:
        """Process PDF files with multiple extraction methods"""
        extracted_text = ""
        
        try:
            # Method 1: PyMuPDF (fitz) - preferred for modern PDFs
            doc = fitz.open(file_path)
            
            pages_text = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                pages_text.append(page_text)
                
                # Extract images from PDF pages if they contain text
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n < 5:  # GRAY or RGB
                            img_data = pix.pil_tobytes(format="PNG")
                            pil_image = Image.open(io.BytesIO(img_data))
                            
                            # OCR on PDF images
                            ocr_result = self._extract_text_from_image(pil_image, f'pdf_page_{page_num}_img_{img_index}')
                            if ocr_result['text'].strip():
                                pages_text.append(f"\n[Image text from page {page_num + 1}]: {ocr_result['text']}")
                        pix = None
                    except Exception as e:
                        logger.warning(f"Error extracting image from PDF page {page_num}: {e}")
            
            extracted_text = '\n\n'.join(pages_text)
            doc.close()
            
            # If PyMuPDF fails or returns little text, try PyPDF2
            if len(extracted_text.strip()) < 100:
                extracted_text = self._extract_with_pypdf2(file_path)
            
        except Exception as e:
            logger.error(f"PyMuPDF extraction failed: {e}")
            # Fallback to PyPDF2
            extracted_text = self._extract_with_pypdf2(file_path)
        
        results['extracted_text'] = extracted_text
        results['document_sections'] = self._extract_pdf_sections(extracted_text)
        
        return results
    
    def _extract_with_pypdf2(self, file_path: str) -> str:
        """Extract text using PyPDF2 as fallback"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                pages_text = []
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    pages_text.append(page_text)
                
                return '\n\n'.join(pages_text)
        except Exception as e:
            logger.error(f"PyPDF2 extraction failed: {e}")
            return ""
    
    def _extract_pdf_sections(self, text: str) -> Dict[str, str]:
        """Extract sections from PDF text"""
        sections = {}
        
        # Simple section detection based on common patterns
        section_patterns = [
            r'^(\d+\.?\s+[A-Z][A-Za-z\s]+)$',
            r'^([A-Z][A-Z\s]{2,})$',
            r'^(CHAPTER\s+\d+[:\-\s]*.*)$',
            r'^(SECTION\s+\d+[:\-\s]*.*)$'
        ]
        
        lines = text.split('\n')
        current_section = "Introduction"
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check if line is a section header
            is_section_header = False
            for pattern in section_patterns:
                if re.match(pattern, line, re.MULTILINE):
                    # Save previous section
                    if current_content:
                        sections[current_section] = '\n'.join(current_content)
                    
                    # Start new section
                    current_section = line
                    current_content = []
                    is_section_header = True
                    break
            
            if not is_section_header:
                current_content.append(line)
        
        # Save last section
        if current_content:
            sections[current_section] = '\n'.join(current_content)
        
        return sections if len(sections) > 1 else {"Full Document": text}
    
    def _process_word_document(self, file_path: str, results: Dict[str, Any]) -> Dict[str, Any]:
        """Process Word documents (DOC, DOCX)"""
        try:
            if file_path.endswith('.docx'):
                doc = DocxDocument(file_path)
                
                # Extract paragraphs
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                
                # Extract tables
                tables_text = []
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append('\t'.join(row_data))
                    tables_text.append('\n'.join(table_data))
                
                # Combine text
                extracted_text = '\n\n'.join(paragraphs)
                if tables_text:
                    extracted_text += '\n\n[TABLES]\n' + '\n\n'.join(tables_text)
                
                results['extracted_text'] = extracted_text
                results['document_sections'] = {
                    'paragraphs': '\n\n'.join(paragraphs),
                    'tables': '\n\n'.join(tables_text) if tables_text else ''
                }
            else:
                # For .doc files, try Tika or other methods
                results = self._process_with_tika(file_path, results)
                
        except Exception as e:
            logger.error(f"Word document processing failed: {e}")
            # Fallback to Tika
            results = self._process_with_tika(file_path, results)
        
        return results
    
    def _process_text_file(self, file_path: str, results: Dict[str, Any]) -> Dict[str, Any]:
        """Process plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise UnicodeDecodeError("Could not decode file with any supported encoding")
            
            results['extracted_text'] = text
            results['document_sections'] = self._extract_text_sections(text)
            
        except Exception as e:
            logger.error(f"Text file processing failed: {e}")
            results['processing_errors'].append(f"Text file error: {str(e)}")
        
        return results
    
    def _extract_text_sections(self, text: str) -> Dict[str, str]:
        """Extract sections from plain text"""
        # Simple section detection
        sections = {}
        lines = text.split('\n')
        
        current_section = "Content"
        current_content = []
        
        for line in lines:
            # Look for section headers (lines that are all caps or start with numbers)
            if (line.strip().isupper() and len(line.strip()) > 3) or \
               re.match(r'^\d+\.?\s+[A-Z]', line.strip()):
                # Save previous section
                if current_content:
                    sections[current_section] = '\n'.join(current_content)
                
                # Start new section
                current_section = line.strip()
                current_content = []
            else:
                current_content.append(line)
        
        # Save last section
        if current_content:
            sections[current_section] = '\n'.join(current_content)
        
        return sections if len(sections) > 1 else {"Full Document": text}
    
    def _process_rtf_file(self, file_path: str, results: Dict[str, Any]) -> Dict[str, Any]:
        """Process RTF files"""
        # RTF processing would require additional libraries
        # For now, fallback to Tika
        return self._process_with_tika(file_path, results)
    
    def _process_presentation(self, file_path: str, results: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process presentation files (PPT, PPTX)"""
        logger.info(f"Processing presentation: {file_path}")
        results['processing_metadata']['processing_methods'].append('presentation_extraction')
        
        try:
            if file_path.endswith('.pptx'):
                prs = Presentation(file_path)
                
                slides_text = []
                slide_details = []
                
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    slide_content = '\n'.join(slide_text)
                    if slide_content:
                        slides_text.append(f"[Slide {i + 1}]\n{slide_content}")
                        slide_details.append({
                            'slide_number': i + 1,
                            'text_content': slide_content,
                            'shape_count': len(slide.shapes)
                        })
                
                results['extracted_text'] = '\n\n'.join(slides_text)
                results['document_sections'] = {
                    'slides': slide_details,
                    'slide_count': len(prs.slides)
                }
            else:
                # For .ppt files, fallback to Tika
                results = self._process_with_tika(file_path, results)
                
        except Exception as e:
            logger.error(f"Presentation processing failed: {e}")
            results['processing_errors'].append(f"Presentation processing error: {str(e)}")
            # Fallback to Tika
            results = self._process_with_tika(file_path, results)
        
        return results
    
    def _process_spreadsheet(self, file_path: str, results: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process spreadsheet files (XLS, XLSX, CSV)"""
        logger.info(f"Processing spreadsheet: {file_path}")
        results['processing_metadata']['processing_methods'].append('spreadsheet_extraction')
        
        try:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                # Read all sheets
                excel_file = pd.ExcelFile(file_path)
                sheets_data = {}
                
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    sheets_data[sheet_name] = df
                
                # Combine all sheets
                all_text = []
                for sheet_name, df in sheets_data.items():
                    sheet_text = self._dataframe_to_text(df, sheet_name)
                    all_text.append(sheet_text)
                
                results['extracted_text'] = '\n\n'.join(all_text)
                results['document_sections'] = {
                    'sheets': {name: self._dataframe_to_text(df, name) for name, df in sheets_data.items()},
                    'sheet_count': len(sheets_data)
                }
                
                return results
            
            # For CSV files
            results['extracted_text'] = self._dataframe_to_text(df, 'Data')
            results['document_sections'] = {
                'data': self._dataframe_to_text(df, 'Data'),
                'rows': len(df),
                'columns': len(df.columns)
            }
            
        except Exception as e:
            logger.error(f"Spreadsheet processing failed: {e}")
            results['processing_errors'].append(f"Spreadsheet processing error: {str(e)}")
        
        return results
    
    def _dataframe_to_text(self, df: pd.DataFrame, sheet_name: str) -> str:
        """Convert DataFrame to readable text"""
        try:
            # Create header
            text_parts = [f"[{sheet_name}]"]
            
            # Add column headers
            headers = ' | '.join(str(col) for col in df.columns)
            text_parts.append(headers)
            text_parts.append('-' * len(headers))
            
            # Add data rows (limit to prevent huge text blocks)
            max_rows = 100  # Limit for performance
            for i, row in df.head(max_rows).iterrows():
                row_text = ' | '.join(str(val) for val in row.values)
                text_parts.append(row_text)
            
            if len(df) > max_rows:
                text_parts.append(f"... ({len(df) - max_rows} more rows)")
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"DataFrame to text conversion failed: {e}")
            return f"[{sheet_name}] - Error processing data"
    
    def _process_with_tika(self, file_path: str, results: Dict[str, Any]) -> Dict[str, Any]:
        """Process file using Apache Tika as fallback"""
        try:
            # Check if Tika server is available
            response = requests.get(f"{TIKA_SERVER_URL}/version", timeout=5)
            if response.status_code != 200:
                raise ServiceUnavailableError("Tika server not available")
            
            # Extract text using Tika
            with open(file_path, 'rb') as file:
                response = requests.put(
                    f"{TIKA_SERVER_URL}/tika",
                    data=file,
                    headers={'Accept': 'text/plain'},
                    timeout=300
                )
            
            if response.status_code == 200:
                results['extracted_text'] = response.text
                results['processing_metadata']['processing_methods'].append('tika_extraction')
            else:
                raise ProcessingError(f"Tika extraction failed with status: {response.status_code}")
                
        except Exception as e:
            logger.error(f"Tika processing failed: {e}")
            results['processing_errors'].append(f"Tika processing error: {str(e)}")
        
        return results
    
    def _process_generic_file(self, file_path: str, results: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process unknown file types"""
        logger.info(f"Processing generic file: {file_path}")
        results['processing_metadata']['processing_methods'].append('generic_processing')
        
        try:
            # Try to read as text file first
            results = self._process_text_file(file_path, results)
            
            # If that fails, try Tika
            if not results['extracted_text'] and not results['processing_errors']:
                results = self._process_with_tika(file_path, results)
                
        except Exception as e:
            logger.error(f"Generic file processing failed: {e}")
            results['processing_errors'].append(f"Generic processing error: {str(e)}")
        
        return results


# Factory function for creating the processor
def create_enhanced_processor() -> EnhancedDocumentProcessor:
    """Create and return an enhanced document processor instance"""
    return EnhancedDocumentProcessor()
