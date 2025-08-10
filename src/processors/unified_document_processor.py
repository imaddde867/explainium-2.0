"""
EXPLAINIUM - Unified Document Processor

A clean, unified processor that handles all document types (images, videos, docs)
and produces structured knowledge output in the expected format.
"""

import os
import io
import json
import logging
import asyncio
from typing import Dict, Any, List, Optional, Tuple, Union
from pathlib import Path
from datetime import datetime
import tempfile
import mimetypes

# Core processing libraries
import requests
import cv2
import numpy as np
from PIL import Image, ImageEnhance
import pandas as pd
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

# OCR and Vision
import pytesseract
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False

# Audio processing
try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False

# Internal imports
from src.logging_config import get_logger
from src.ai.intelligent_knowledge_extractor import IntelligentKnowledgeExtractor
from src.exceptions import ProcessingError, AIError

logger = get_logger(__name__)


class UnifiedDocumentProcessor:
    """
    Unified document processor that handles all document types and produces
    structured knowledge output in the expected format.
    """
    
    def __init__(self):
        self.knowledge_extractor = IntelligentKnowledgeExtractor()
        self.supported_formats = {
            'text': ['.pdf', '.doc', '.docx', '.txt', '.rtf'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'],
            'spreadsheet': ['.xls', '.xlsx', '.csv'],
            'presentation': ['.ppt', '.pptx'],
            'audio': ['.mp3', '.wav', '.flac', '.aac'],
            'video': ['.mp4', '.avi', '.mov', '.mkv']
        }
        
        # Initialize processing capabilities
        self._init_ocr()
        self._init_audio_processing()
        
        logger.info("Unified Document Processor initialized")
    
    def _init_ocr(self):
        """Initialize OCR capabilities"""
        try:
            # Test pytesseract
            pytesseract.get_tesseract_version()
            self.ocr_available = True
            logger.info("OCR initialized successfully")
        except Exception as e:
            logger.warning(f"OCR initialization failed: {e}")
            self.ocr_available = False
    
    def _init_audio_processing(self):
        """Initialize audio processing capabilities"""
        try:
            if WHISPER_AVAILABLE:
                self.whisper_model = whisper.load_model("base")
                self.audio_processing_available = True
                logger.info("Audio processing initialized successfully")
            else:
                self.audio_processing_available = False
                self.whisper_model = None
        except Exception as e:
            logger.warning(f"Audio processing initialization failed: {e}")
            self.audio_processing_available = False
            self.whisper_model = None
    
    async def process_document(self, file_path: str, document_id: int) -> Dict[str, Any]:
        """
        Main method to process any supported document format
        
        Args:
            file_path: Path to the document file
            document_id: Database ID of the document
            
        Returns:
            Dictionary containing extracted content and knowledge
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise ProcessingError(f"File not found: {file_path}")
        
        file_extension = file_path.suffix.lower()
        file_type = self._get_file_type(file_extension)
        
        logger.info(f"Processing {file_type} document: {file_path.name}")
        
        try:
            # Extract content based on file type
            if file_type == 'text':
                content = await self._process_text_document(file_path)
            elif file_type == 'image':
                content = await self._process_image_document(file_path)
            elif file_type == 'spreadsheet':
                content = await self._process_spreadsheet_document(file_path)
            elif file_type == 'presentation':
                content = await self._process_presentation_document(file_path)
            elif file_type == 'audio':
                content = await self._process_audio_document(file_path)
            elif file_type == 'video':
                content = await self._process_video_document(file_path)
            else:
                raise ProcessingError(f"Unsupported file type: {file_extension}")
            
            # Extract knowledge from content using intelligent extractor
            knowledge = await self.knowledge_extractor.extract_knowledge(
                content.get('text', ''),
                file_type
            )
            
            result = {
                'document_id': document_id,
                'filename': file_path.name,
                'file_type': file_type,
                'processed_at': datetime.now().isoformat(),
                'content': content,
                'knowledge': knowledge,
                'processing_status': 'completed'
            }
            
            logger.info(f"Successfully processed document: {file_path.name}")
            return result
            
        except Exception as e:
            logger.error(f"Error processing document {file_path.name}: {e}")
            raise ProcessingError(f"Failed to process document: {str(e)}")
    
    def _get_file_type(self, extension: str) -> str:
        """Determine file type from extension"""
        for file_type, extensions in self.supported_formats.items():
            if extension in extensions:
                return file_type
        return 'unknown'
    
    async def _process_text_document(self, file_path: Path) -> Dict[str, Any]:
        """Process text-based documents (PDF, DOC, DOCX, TXT)"""
        extension = file_path.suffix.lower()
        
        if extension == '.pdf':
            return await self._extract_pdf_content(file_path)
        elif extension in ['.doc', '.docx']:
            return await self._extract_word_content(file_path)
        elif extension == '.txt':
            return await self._extract_text_content(file_path)
        else:
            # Try Tika for other formats
            return await self._extract_with_tika(file_path)
    
    async def _extract_pdf_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PDF files"""
        text_content = ""
        metadata = {}
        
        try:
            # Try PyMuPDF first (better for complex PDFs)
            doc = fitz.open(file_path)
            page_count = len(doc)
            
            for page_num in range(page_count):
                page = doc.load_page(page_num)
                text_content += page.get_text()
            
            metadata = doc.metadata
            doc.close()
            
        except Exception as e:
            logger.warning(f"PyMuPDF failed, trying PyPDF2: {e}")
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    page_count = len(pdf_reader.pages)
                    
                    for page in pdf_reader.pages:
                        text_content += page.extract_text()
                    
                    metadata = pdf_reader.metadata
                    
            except Exception as e2:
                logger.error(f"Both PDF readers failed: {e2}")
                raise ProcessingError(f"Failed to extract PDF content: {str(e2)}")
        
        return {
            'text': text_content,
            'page_count': page_count,
            'metadata': metadata,
            'language': 'en',  # Default assumption
            'extraction_method': 'pdf_reader'
        }
    
    async def _extract_word_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Word documents"""
        try:
            doc = DocxDocument(file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            return {
                'text': text_content,
                'page_count': None,  # Word doesn't have fixed page count
                'metadata': {
                    'author': doc.core_properties.author,
                    'title': doc.core_properties.title,
                    'subject': doc.core_properties.subject,
                    'created': doc.core_properties.created,
                    'modified': doc.core_properties.modified
                },
                'language': 'en',
                'extraction_method': 'python_docx'
            }
            
        except Exception as e:
            logger.error(f"Failed to extract Word content: {e}")
            raise ProcessingError(f"Failed to extract Word content: {str(e)}")
    
    async def _extract_text_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text_content = file.read()
            
            return {
                'text': text_content,
                'page_count': 1,
                'metadata': {},
                'language': 'en',
                'extraction_method': 'direct_read'
            }
            
        except Exception as e:
            logger.error(f"Failed to extract text content: {e}")
            raise ProcessingError(f"Failed to extract text content: {str(e)}")
    
    async def _process_image_document(self, file_path: Path) -> Dict[str, Any]:
        """Process image documents with OCR"""
        try:
            # Load and enhance image
            image = Image.open(file_path)
            enhanced_image = self._enhance_image_for_ocr(image)
            
            # Extract text using OCR
            if self.ocr_available:
                text_content = pytesseract.image_to_string(enhanced_image)
                
                # Try EasyOCR if available and pytesseract fails
                if not text_content.strip() and EASYOCR_AVAILABLE:
                    reader = easyocr.Reader(['en'])
                    results = reader.readtext(np.array(enhanced_image))
                    text_content = ' '.join([result[1] for result in results])
                
                return {
                    'text': text_content,
                    'page_count': 1,
                    'metadata': {
                        'image_size': image.size,
                        'image_mode': image.mode,
                        'ocr_confidence': self._calculate_ocr_confidence(text_content)
                    },
                    'language': 'en',
                    'extraction_method': 'ocr'
                }
            else:
                raise ProcessingError("OCR not available for image processing")
                
        except Exception as e:
            logger.error(f"Failed to process image: {e}")
            raise ProcessingError(f"Failed to process image: {str(e)}")
    
    def _enhance_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Enhance image for better OCR results"""
        try:
            # Convert to grayscale if not already
            if image.mode != 'L':
                image = image.convert('L')
            
            # Enhance contrast
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(2.0)
            
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(2.0)
            
            return image
            
        except Exception as e:
            logger.warning(f"Image enhancement failed: {e}")
            return image
    
    async def _process_spreadsheet_document(self, file_path: Path) -> Dict[str, Any]:
        """Process spreadsheet documents"""
        try:
            extension = file_path.suffix.lower()
            
            if extension == '.csv':
                df = pd.read_csv(file_path)
            elif extension in ['.xls', '.xlsx']:
                df = pd.read_excel(file_path, sheet_name=None)
                # Combine all sheets
                if isinstance(df, dict):
                    all_sheets = []
                    for sheet_name, sheet_df in df.items():
                        all_sheets.append(f"=== {sheet_name} ===")
                        all_sheets.append(self._dataframe_to_text(sheet_df))
                    text_content = '\n'.join(all_sheets)
                else:
                    text_content = self._dataframe_to_text(df)
            else:
                raise ProcessingError(f"Unsupported spreadsheet format: {extension}")
            
            return {
                'text': text_content,
                'page_count': 1,
                'metadata': {
                    'rows': len(df) if not isinstance(df, dict) else sum(len(sheet) for sheet in df.values()),
                    'columns': len(df.columns) if not isinstance(df, dict) else max(len(sheet.columns) for sheet in df.values())
                },
                'language': 'en',
                'extraction_method': 'pandas'
            }
            
        except Exception as e:
            logger.error(f"Failed to process spreadsheet: {e}")
            raise ProcessingError(f"Failed to process spreadsheet: {str(e)}")
    
    def _dataframe_to_text(self, df: pd.DataFrame) -> str:
        """Convert DataFrame to readable text"""
        try:
            # Convert to string representation
            text_content = df.to_string(index=False)
            return text_content
        except Exception as e:
            logger.warning(f"DataFrame to text conversion failed: {e}")
            return str(df)
    
    async def _process_presentation_document(self, file_path: Path) -> Dict[str, Any]:
        """Process presentation documents"""
        try:
            prs = Presentation(file_path)
            text_content = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content += f"\n=== Slide {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
                
                text_content += "\n"
            
            return {
                'text': text_content,
                'page_count': len(prs.slides),
                'metadata': {
                    'slide_count': len(prs.slides),
                    'presentation_format': 'PowerPoint'
                },
                'language': 'en',
                'extraction_method': 'python_pptx'
            }
            
        except Exception as e:
            logger.error(f"Failed to process presentation: {e}")
            raise ProcessingError(f"Failed to process presentation: {str(e)}")
    
    async def _process_audio_document(self, file_path: Path) -> Dict[str, Any]:
        """Process audio documents with transcription"""
        if not self.audio_processing_available:
            raise ProcessingError("Audio processing not available")
        
        try:
            # Transcribe audio using Whisper
            result = self.whisper_model.transcribe(str(file_path))
            text_content = result["text"]
            
            return {
                'text': text_content,
                'page_count': 1,
                'metadata': {
                    'audio_duration': result.get("duration"),
                    'language': result.get("language", "en"),
                    'transcription_confidence': result.get("confidence", 0.0)
                },
                'language': result.get("language", "en"),
                'extraction_method': 'whisper'
            }
            
        except Exception as e:
            logger.error(f"Failed to process audio: {e}")
            raise ProcessingError(f"Failed to process audio: {str(e)}")
    
    async def _process_video_document(self, file_path: Path) -> Dict[str, Any]:
        """Process video documents"""
        try:
            # Extract audio from video
            temp_audio = tempfile.NamedTemporaryFile(suffix='.wav', delete=False)
            temp_audio_path = temp_audio.name
            temp_audio.close()
            
            # Use ffmpeg to extract audio
            import subprocess
            try:
                subprocess.run([
                    'ffmpeg', '-i', str(file_path), '-vn', '-acodec', 'pcm_s16le',
                    '-ar', '16000', '-ac', '1', temp_audio_path
                ], check=True, capture_output=True)
                
                # Process the extracted audio
                audio_content = await self._process_audio_document(Path(temp_audio_path))
                
                # Get video metadata
                video_info = self._get_video_info(file_path)
                
                # Clean up temp file
                os.unlink(temp_audio_path)
                
                return {
                    'text': audio_content['text'],
                    'page_count': 1,
                    'metadata': {
                        **audio_content['metadata'],
                        **video_info
                    },
                    'language': audio_content['language'],
                    'extraction_method': 'video_audio_extraction'
                }
                
            except subprocess.CalledProcessError as e:
                raise ProcessingError(f"Failed to extract audio from video: {str(e)}")
            finally:
                # Ensure temp file is cleaned up
                if os.path.exists(temp_audio_path):
                    os.unlink(temp_audio_path)
                    
        except Exception as e:
            logger.error(f"Failed to process video: {e}")
            raise ProcessingError(f"Failed to process video: {str(e)}")
    
    def _get_video_info(self, video_path: Path) -> Dict[str, Any]:
        """Get basic video information"""
        try:
            cap = cv2.VideoCapture(str(video_path))
            if cap.isOpened():
                fps = cap.get(cv2.CAP_PROP_FPS)
                frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                duration = frame_count / fps if fps > 0 else 0
                width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                cap.release()
                
                return {
                    'video_duration': duration,
                    'video_fps': fps,
                    'video_frames': frame_count,
                    'video_resolution': f"{width}x{height}"
                }
            else:
                return {}
        except Exception as e:
            logger.warning(f"Failed to get video info: {e}")
            return {}
    
    async def _extract_with_tika(self, file_path: Path) -> Dict[str, Any]:
        """Extract content using Apache Tika (fallback method)"""
        try:
            # This would integrate with Tika server
            # For now, return a basic error
            raise ProcessingError("Tika extraction not yet implemented")
        except Exception as e:
            logger.error(f"Tika extraction failed: {e}")
            raise ProcessingError(f"Tika extraction failed: {str(e)}")
    
    def _calculate_ocr_confidence(self, text: str) -> float:
        """Calculate confidence score for OCR results"""
        if not text.strip():
            return 0.0
        
        # Simple heuristic: longer text with fewer special characters = higher confidence
        text_length = len(text)
        special_chars = len(re.findall(r'[^a-zA-Z0-9\s\.\,\;\:\!\?]', text))
        
        if text_length == 0:
            return 0.0
        
        # Base confidence on text quality
        base_confidence = min(0.9, text_length / 1000)  # More text = higher confidence
        
        # Penalize for special characters
        char_penalty = min(0.3, special_chars / text_length)
        
        return max(0.1, base_confidence - char_penalty)
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get supported file formats"""
        return self.supported_formats
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        return {
            'ocr_available': self.ocr_available,
            'audio_processing_available': self.audio_processing_available,
            'supported_formats': len(self.supported_formats),
            'total_supported_extensions': sum(len(exts) for exts in self.supported_formats.values())
        }